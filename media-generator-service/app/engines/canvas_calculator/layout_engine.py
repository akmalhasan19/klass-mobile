"""Canvas layout engine — geometry + orchestration for the PPTX fallback.

When a ``SlideBlueprint`` slide exceeds the master template's capacity, this
engine computes a grid of card boxes (plus a title band) and delegates the
actual drawing of each box to an injected ``ShapeRenderer`` (implemented in
``shape_renderer.py``, Task 3.7). It satisfies the ``CanvasRenderer`` Protocol
that ``TemplateInjector`` calls via ``canvas_engine.render_slide(...)``.

Deviation from the plan's prototype (§5b)
-----------------------------------------
The prototype renders rounded-rectangle shapes inline and uses *fixed* grid
cell heights (``cell_h = usable_h // rows``). Two improvements here:

1. **Drawing is delegated**, not embedded. Task 3.6 is the calculator; the
   shape drawing is Task 3.7. Coupling them via a ``ShapeRenderer`` Protocol
   (mirroring the ``CanvasRenderer`` Protocol already used by the injector)
   keeps the two tasks cleanly separated and lets the engine be unit-tested
   with a stub renderer today.

2. **Card heights are fit-aware.** We reuse ``font_metrics.estimate_box``
   (Task 3.5) to size each card box so its text actually fits, instead of
   clipping overflowing content into a fixed-height cell. Column widths and
   x-positions still come from the grid; only the per-row height adapts to the
   tallest card in that row.
"""
from __future__ import annotations

import math
from dataclasses import dataclass
from typing import Protocol, runtime_checkable

from pptx import Presentation
from pptx.util import Emu, Inches

from app.engines.blueprint import Card, ContentBlock, Slide
from app.engines.canvas_calculator.font_metrics import estimate_box
from app.engines.canvas_calculator.shape_renderer import _font_sizes
from app.errors import GenerationError


@dataclass(frozen=True)
class Box:
    """A rectangle in EMU (origin top-left)."""

    x: Emu
    y: Emu
    w: Emu
    h: Emu


@dataclass(frozen=True)
class SlideLayout:
    """Computed geometry for one canvas-rendered slide."""

    title_box: Box
    card_boxes: list[Box]  # aligned 1:1 with ``slide.cards``


@runtime_checkable
class ShapeRenderer(Protocol):
    """Structural type for the concrete shape renderer (Task 3.7)."""

    def draw_title(self, slide, title: str, box: Box) -> None:
        """Draw *title* into *box* on *slide*."""
        ...

    def draw_card(self, slide, card: Card, box: Box) -> None:
        """Draw *card* into *box* on *slide*."""
        ...


class CanvasLayoutEngine:
    """Compute canvas slide geometry and delegate shape drawing.

    Implements the ``CanvasRenderer`` Protocol consumed by
    ``TemplateInjector`` (``render_slide(presentation, slide)``).
    """

    def __init__(
        self,
        slide_width: Emu = Inches(13.333),
        slide_height: Emu = Inches(7.5),
        margin: Emu = Inches(0.5),
        gap: Emu = Inches(0.25),
        title_band: Emu = Inches(1.1),
        renderer: ShapeRenderer | None = None,
    ) -> None:
        self._w, self._h = slide_width, slide_height
        self._margin = margin
        self._gap = gap
        self._title_band = title_band
        self._renderer = renderer

    def render_slide(self, presentation: Presentation, slide: Slide) -> SlideLayout:
        """Append a canvas-calculated slide for *slide* to *presentation*.

        Raises:
            GenerationError: If no ``ShapeRenderer`` was injected (the concrete
                renderer from Task 3.7 must be supplied before use).
        """
        if self._renderer is None:
            raise GenerationError(
                code="canvas_renderer_missing",
                message=(
                    "CanvasLayoutEngine has no ShapeRenderer injected. "
                    "Construct it with renderer=ShapeRenderer() (app.engines."
                    "canvas_calculator.shape_renderer) before rendering."
                ),
            )

        layout = self.compute_layout(slide)
        prs_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._renderer.draw_title(prs_slide, slide.title, layout.title_box)
        for card, box in zip(slide.cards, layout.card_boxes):
            self._renderer.draw_card(prs_slide, card, box)
        return layout

    def compute_layout(self, slide: Slide) -> SlideLayout:
        """Compute title band + per-card boxes, fit-aware to card content.

        Column count follows the plan's ``_pick_columns`` heuristic; within
        each row the height adapts to the tallest card so text never clips.
        """
        title_box = Box(
            self._margin,
            self._margin,
            self._w - 2 * self._margin,
            self._title_band,
        )

        columns = self._pick_columns(len(slide.cards))
        rows = math.ceil(len(slide.cards) / columns)

        usable_w = self._w - 2 * self._margin - (columns - 1) * self._gap
        usable_h = (
            self._h
            - self._margin
            - self._title_band
            - self._margin
            - (rows - 1) * self._gap
        )
        cell_w = usable_w // columns
        min_cell_h = usable_h // rows

        card_boxes: list[Box] = []
        y = self._margin + self._title_band
        for r in range(rows):
            row_cards = slide.cards[r * columns : (r + 1) * columns]
            row_h = max(
                min_cell_h,
                max((self._estimate_card_height(c, cell_w) for c in row_cards), default=min_cell_h),
            )
            x = self._margin
            for _ in row_cards:
                card_boxes.append(Box(x, y, cell_w, row_h))
                x += cell_w + self._gap
            y += row_h + self._gap

        # Guard against pathological overflow: clamp the total to slide height.
        if y - self._gap > self._h:
            raise GenerationError(
                code="canvas_layout_overflow",
                message=(
                    f"Canvas layout for slide '{slide.title}' exceeds slide "
                    f"height ({int(y - self._gap)} > {int(self._h)} EMU); "
                    f"reduce card count."
                ),
            )

        return SlideLayout(title_box, card_boxes)

    @staticmethod
    def _pick_columns(card_count: int) -> int:
        # 1→1, 2→2, 3-4→2, 5-9→3, ≥10→4 (plan heuristic)
        if card_count <= 1:
            return 1
        if card_count <= 4:
            return 2
        if card_count <= 9:
            return 3
        return 4

    def _estimate_card_height(self, card: Card, box_w_emu: Emu) -> Emu:
        """Estimate a card's required height from its text via ``font_metrics``.

        Uses the same per-width font sizes (``_font_sizes``) that
        ``CanvasShapeRenderer`` applies when drawing, so the computed box height
        matches the rendered text and never clips.

        The available text width is ``box_w_emu - 2 * text_frame_h_margin``
        because the ``CanvasShapeRenderer`` sets
        ``frame.margin_left = frame.margin_right = 0.12 in`` on every card
        shape.  Failing to subtract this caused ``estimate_box`` to
        underestimate the number of wrapped lines, which in turn produced
        card boxes that were too short — the root cause of the content
        overflow the user observed.

        Paragraph spacing (``space_after``) is also accounted for: the
        heading paragraph has ``Pt(6)`` and each body paragraph has
        ``Pt(4)``.
        """
        # Text frame horizontal margins (must match shape_renderer.py)
        text_h_margin = Inches(0.12) * 2
        text_w_emu = box_w_emu - int(text_h_margin)
        if text_w_emu <= 0:
            text_w_emu = box_w_emu  # fallback for extremely narrow boxes

        # Vertical margins inside the text frame (must match shape_renderer.py)
        text_v_margin = Inches(0.1) * 2  # top + bottom

        heading_pt, body_pt = _font_sizes(box_w_emu)
        total: Emu = Emu(0)
        if card.heading:
            total += estimate_box(card.heading, heading_pt, text_w_emu)
            total += Pt(6)  # paragraph.space_after for heading
        body_blocks = card.body_blocks
        body_text = "\n".join(self._format_block(b) for b in body_blocks)
        if body_text:
            total += estimate_box(body_text, body_pt, text_w_emu)
            # Account for space_after on each body paragraph
            total += Pt(4) * len(body_blocks)
        return total + text_v_margin

    @staticmethod
    def _format_block(block: ContentBlock) -> str:
        if block.kind == "bullet":
            return f"• {block.content}"
        if block.kind == "checklist":
            return f"☐ {block.content}"
        if block.kind == "note":
            return f"Note: {block.content}"
        return block.content

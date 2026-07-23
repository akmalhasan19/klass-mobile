"""Concrete shape renderer for the Canvas Calculator fallback (Task 3.7).

Implements the :class:`~app.engines.canvas_calculator.layout_engine.ShapeRenderer`
Protocol. The layout engine computes ``Box`` rectangles (in EMU); this module
turns each box into real ``python-pptx`` shapes:

* ``draw_title`` — a text box for the slide title band.
* ``draw_card``  — a rounded-rectangle card with a heading (bold) and the card's
  body blocks rendered as bullets / checklists / paragraphs / notes.

Deviation from the plan's prototype (§5b)
-----------------------------------------
The prototype inlines the drawing inside ``layout_engine`` and uses *fixed*
``14pt / 11pt`` fonts. This implementation keeps drawing in its own module
(per the plan's file split) and scales the font to the box width, honouring the
plan's "font size per column count" requirement without hard-coding a column
count that the renderer never receives. Because ``CanvasLayoutEngine`` derives
column count from the same box width, the two stay in lock-step.

Styling follows the Klass educational palette already used in
``app/generators/pptx_generator.py`` (and mirrored in the Marp theme CSS), so
canvas-fallback slides stay visually consistent with template-injected slides
and the HTML preview. ``ContentBlock`` kinds are the four real ones
(``paragraph``, ``bullet``, ``checklist``, ``note``), matching the blueprint,
not the prototype's three.
"""
from __future__ import annotations

from typing import TYPE_CHECKING

# pyrefly: ignore [missing-import]
from pptx.dml.color import RGBColor
# pyrefly: ignore [missing-import]
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
# pyrefly: ignore [missing-import]
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
# pyrefly: ignore [missing-import]
from pptx.util import Inches, Pt

from app.engines.blueprint import Card, ContentBlock

if TYPE_CHECKING:  # avoid runtime circular import with layout_engine.Box
    from app.engines.canvas_calculator.layout_engine import Box


# --- Klass educational palette (kept in sync with pptx_generator.py) ----------
_PRIMARY = RGBColor(11, 31, 51)     # title text        #0B1F33
_BODY = RGBColor(31, 41, 51)        # body text         #1F2933
_MUTED = RGBColor(82, 96, 109)      # subtitle / hint   #52606D
_ACCENT = RGBColor(15, 76, 92)      # heading accent    #0F4C5C
_BORDER = RGBColor(188, 204, 220)   # card / box border #BCCCDC
_CARD_FILL = RGBColor(255, 255, 255)

_EMU_PER_INCH = 914_400

# About 8% of the shorter side — a soft, modern corner on the cards.
_CORNER_RADIUS = 0.08


def _font_sizes(width_emu: int) -> tuple[int, int]:
    """Return ``(heading_pt, body_pt)`` scaled to the card *width_emu*.

    Wider boxes (fewer columns) get larger type; narrow boxes (more columns)
    shrink so text stays legible without clipping. ``CanvasLayoutEngine``
    already derives the column count from this same width, so the estimate
    (``font_metrics``) and the actual render use identical sizes.
    """
    width_in = width_emu / _EMU_PER_INCH
    if width_in >= 6.0:      # 1 column
        return 18, 13
    if width_in >= 4.0:      # 2 columns
        return 15, 12
    if width_in >= 2.8:      # 3 columns
        return 13, 11
    return 12, 10            # 4 (or more) columns


class CanvasShapeRenderer:
    """Draws canvas-fallback shapes (title band + cards) into a PPTX slide.

    Satisfies the ``ShapeRenderer`` Protocol consumed by
    ``CanvasLayoutEngine.render_slide``.
    """

    def draw_title(self, slide, title: str, box: "Box") -> None:
        """Draw *title* into the title band *box* on *slide*."""
        textbox = slide.shapes.add_textbox(box.x, box.y, box.w, box.h)
        frame = textbox.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        frame.margin_left = frame.margin_right = Inches(0.1)
        frame.margin_top = frame.margin_bottom = Inches(0.05)

        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = title
        run.font.bold = True
        run.font.size = Pt(24)
        run.font.color.rgb = _PRIMARY
        run.font.name = "Calibri"

    def draw_card(self, slide, card: Card, box: "Box") -> None:
        """Draw *card* into the card *box* on *slide*."""
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, box.x, box.y, box.w, box.h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _CARD_FILL
        shape.line.color.rgb = _BORDER
        shape.line.width = Pt(1)
        self._set_corner_radius(shape, _CORNER_RADIUS)

        frame = shape.text_frame
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = Inches(0.12)
        frame.margin_top = frame.margin_bottom = Inches(0.1)
        # ``add_shape`` seeds a default "Rounded Rectangle" run — clear it so
        # our own content starts from a clean paragraph.
        frame.text = ""

        heading_pt, body_pt = _font_sizes(box.w)

        first = True
        if card.heading:
            paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
            first = False
            self._add_run(
                paragraph, card.heading, bold=True, size=heading_pt, color=_ACCENT
            )
            paragraph.space_after = Pt(6)

        for block in card.body_blocks:
            paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
            first = False
            self._add_run(
                paragraph,
                self._format_block(block),
                bold=False,
                size=body_pt,
                color=_BODY,
            )
            paragraph.space_after = Pt(4)

    # ------------------------------------------------------------------ helpers
    @staticmethod
    def _add_run(paragraph, text: str, *, bold: bool, size: int, color: RGBColor):
        run = paragraph.add_run()
        run.text = text
        run.font.bold = bold
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        return run

    @staticmethod
    def _format_block(block: ContentBlock) -> str:
        """Render a body block's prefix, matching the injector's convention."""
        if block.kind == "bullet":
            return f"• {block.content}"
        if block.kind == "checklist":
            return f"☐ {block.content}"
        if block.kind == "note":
            return f"Note: {block.content}"
        return block.content

    @staticmethod
    def _set_corner_radius(shape, fraction: float) -> None:
        """Set the rounded-rectangle corner radius (fraction of shorter side)."""
        try:
            shape.adjustments[0] = fraction
        except Exception:
            # Not all builds expose adjustments uniformly; ignore if unavailable.
            pass

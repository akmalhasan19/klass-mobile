"""Template injector: fill a master ``.pptx`` with ``SlideBlueprint`` content.

Deviation from the plan's prototype (§5a)
-----------------------------------------
The plan's prototype calls ``presentation.slides.add_slide(layout.slide_layout)``
to obtain a target slide. That is **incompatible** with the master template
actually shipped in ``app/templates/masters/klass-educational-v1.pptx``.

Per ``app/templates/registry.py`` (the master contract), ``slide_index`` indexes
``Presentation.slides`` — the master ships three *fully designed slides*
(title / content / assessment) whose placeholder shapes live on those slides,
not on bare slide layouts. ``add_slide(slide_layout)`` would therefore yield an
empty layout without the designer's shapes. This injector instead **duplicates
the source slide** (``master.slides[slide_index]``) into the output deck and
fills its shapes by ``shape.name``.

Other corrections vs the prototype:
- ``ContentBlock.kind`` is ``"paragraph" | "bullet" | "checklist" | "note"``
  (not just bullet/checklist), so ``_format_block`` maps all four kinds.
- Multi-line body text preserves the master's run formatting: the first line
  reuses the existing first run; subsequent lines replicate that run's font so
  bullet/colour styling stays consistent across the whole body.
- ``notes`` is written to the slide's notes pane, not resolved as a shape.
- Overflow / no-layout slides are delegated to an optional ``canvas_engine``
  (the Canvas Calculator, implemented in a later task); when none is supplied
  they are recorded in ``fallback_slides`` for the orchestrator to render.
"""
from __future__ import annotations

from copy import deepcopy
from dataclasses import dataclass, field
from pathlib import Path
from typing import Protocol, runtime_checkable

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.shapes.base import BaseShape
from pptx.slide import Slide as PptxSlide

from app.engines.blueprint import ContentBlock, Slide, SlideBlueprint
from app.engines.pptx_injector.manifest import LayoutManifest, TemplateManifest
from app.engines.pptx_injector.placeholder_resolver import resolve_shape


@dataclass(frozen=True)
class InjectionResult:
    """Outcome of a template injection pass.

    Attributes:
        slide_count: Number of slides actually written to the output deck
            (template-filled slides; canvas slides are added by the
            orchestrator, not here).
        fallback_slides: 1-indexed slide numbers that could not be template
            filled (capacity overflow or no matching layout) and must be
            rendered by the Canvas Calculator.
        warnings: Non-fatal issues encountered while injecting.
    """

    slide_count: int
    fallback_slides: list[int] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)


@runtime_checkable
class CanvasRenderer(Protocol):
    """Structural type for the Canvas Calculator fallback engine.

    Declared as a Protocol (rather than importing the not-yet-built
    ``canvas_calculator`` module) so this injector can be used and tested
    without the fallback engine existing yet.
    """

    def render_slide(self, presentation: Presentation, slide: Slide) -> object:
        """Append a canvas-calculated slide for *slide* to *presentation*."""
        ...


def _copy_run_formatting(src_font, dst_font) -> None:
    """Replicate run font attributes from *src_font* onto *dst_font*."""
    dst_font.bold = src_font.bold
    dst_font.italic = src_font.italic
    dst_font.size = src_font.size
    dst_font.name = src_font.name
    if src_font.color is not None and src_font.color.type is not None:
        try:
            dst_font.color.rgb = src_font.color.rgb
        except Exception:
            pass


class TemplateInjector:
    """Inject ``SlideBlueprint`` content into a master ``.pptx`` template."""

    def __init__(
        self,
        master_path: Path,
        manifest: TemplateManifest,
        canvas_engine: CanvasRenderer | None = None,
    ) -> None:
        self._master_path = master_path
        self._manifest = manifest
        self._canvas_engine = canvas_engine

    def inject(self, blueprint: SlideBlueprint, output_path: Path) -> InjectionResult:
        """Fill *blueprint* into a copy of the master and save to *output_path*.

        Opens the master ``.pptx`` per request (``python-pptx`` is not
        thread-safe), duplicates the designed source slide for each
        template-fitting blueprint slide, and records overflow / unmatched
        slides in ``fallback_slides``.
        """
        master_prs = Presentation(str(self._master_path))
        output = Presentation()
        output.slide_width = master_prs.slide_width
        output.slide_height = master_prs.slide_height

        warnings: list[str] = []
        fallback: list[int] = []

        for index, slide in enumerate(blueprint.slides, start=1):
            layout = self._manifest.pick_layout(slide.slide_type)
            if layout is None:
                warnings.append(
                    f"slide {index}: no layout for '{slide.slide_type}' "
                    f"-> canvas fallback"
                )
                self._delegate_canvas(index, output, slide, fallback)
                continue

            if self._fits(layout, slide):
                self._fill_slide(output, master_prs, layout, slide)
            else:
                warnings.append(
                    f"slide {index}: exceeds '{layout.layout_id}' capacity "
                    f"-> canvas fallback"
                )
                self._delegate_canvas(index, output, slide, fallback)

        output.save(str(output_path))
        return InjectionResult(
            slide_count=len(output.slides),
            fallback_slides=fallback,
            warnings=warnings,
        )

    def _delegate_canvas(
        self, index: int, output: Presentation, slide: Slide, fallback: list[int]
    ) -> None:
        if self._canvas_engine is not None:
            self._canvas_engine.render_slide(output, slide)
            slide.layout_source = "canvas"
        fallback.append(index)

    def _fits(self, layout: LayoutManifest, slide: Slide) -> bool:
        body = layout.placeholder("body")
        if body is None or body.capacity is None:
            return True
        cap = body.capacity
        if cap.max_cards and len(slide.cards) > cap.max_cards:
            return False
        if cap.max_chars:
            total = sum(
                len(card.heading or "")
                + sum(len(block.content) for block in card.body_blocks)
                for card in slide.cards
            )
            if total > cap.max_chars:
                return False
        return True

    def _fill_slide(
        self,
        output: Presentation,
        master_prs: Presentation,
        layout: LayoutManifest,
        slide: Slide,
    ) -> None:
        source = master_prs.slides[layout.slide_index]
        prs_slide = self._duplicate_slide(output, source)

        for placeholder_id, value in self._bindings(slide).items():
            if not value:
                continue
            # Notes are written to the slide's notes pane, not to a
            # placeholder shape — handle them before the spec lookup so
            # the injector does not skip them when the manifest has no
            # "notes" placeholder entry.
            if placeholder_id == "notes":
                self._write_notes(prs_slide, value)
                continue
            spec = layout.placeholder(placeholder_id)
            if spec is None:
                continue
            if spec.kind != "text":
                continue
            shape = resolve_shape(prs_slide, spec.shape_name, warnings=[])
            if shape is None:
                continue
            self._write_text(shape, value)

        slide.layout_source = "template"

    def _duplicate_slide(self, output: Presentation, source: PptxSlide) -> PptxSlide:
        """Duplicate a designed *source* slide into *output*.

        The master ships fully-designed slides, so we copy the slide's shape
        tree (and background) rather than instantiating a bare layout. A blank
        layout is used only as the container; its default shapes are stripped
        before the source shapes are appended.
        """
        new_slide = output.slides.add_slide(output.slide_layouts[6])
        sp_tree = new_slide.shapes._spTree
        for shape in list(new_slide.shapes):
            sp_tree.remove(shape._element)

        src_sp_tree = source.shapes._spTree
        for child in src_sp_tree:
            sp_tree.append(deepcopy(child))

        self._copy_slide_background(new_slide, source)
        return new_slide

    @staticmethod
    def _copy_slide_background(dst_slide: PptxSlide, src_slide: PptxSlide) -> None:
        src_cSld = src_slide._element.find(qn("p:cSld"))
        if src_cSld is None:
            return
        src_bg = src_cSld.find(qn("p:bg"))
        if src_bg is None:
            return
        dst_cSld = dst_slide._element.find(qn("p:cSld"))
        if dst_cSld is None:
            return
        dst_bg = dst_cSld.find(qn("p:bg"))
        if dst_bg is not None:
            dst_cSld.remove(dst_bg)
        dst_cSld.append(deepcopy(src_bg))

    def _bindings(self, slide: Slide) -> dict[str, str]:
        out: dict[str, str] = {"title": slide.title}
        if slide.subtitle:
            out["subtitle"] = slide.subtitle
        lines: list[str] = []
        for card in slide.cards:
            if card.heading:
                lines.append(card.heading)
            lines.extend(self._format_block(b) for b in card.body_blocks)
        out["body"] = "\n".join(lines)
        if slide.speaker_notes:
            out["notes"] = slide.speaker_notes
        return out

    @staticmethod
    def _format_block(block: ContentBlock) -> str:
        if block.kind == "bullet":
            return f"• {block.content}"
        if block.kind == "checklist":
            return f"☐ {block.content}"
        if block.kind == "note":
            return f"Note: {block.content}"
        return block.content

    @staticmethod
    def _write_notes(prs_slide: PptxSlide, value: str) -> None:
        frame = prs_slide.notes_slide.notes_text_frame
        frame.text = value

    @staticmethod
    def _write_text(shape: BaseShape, value: str, *, preserve_master: bool = True) -> None:
        """Write *value* into *shape*'s text frame, preserving master formatting.

        For single-line placeholders (title/subtitle) the master's first run is
        reused verbatim. For multi-line bodies the first line reuses the master's
        first run and each subsequent line gets a new paragraph whose run
        replicates the master run's font, so bullet/colour styling is consistent.

        **Overflow protection**: When the rendered text would exceed the
        shape's available height, the font size is iteratively reduced
        (in 0.5 pt steps, down to a floor of 8 pt) until the content
        fits.  This prevents the user-observed issue where long slide
        content overflows the PPT view and appears truncated.
        """
        frame = shape.text_frame
        if not preserve_master or not frame.paragraphs or not frame.paragraphs[0].runs:
            frame.text = value
            return

        lines = value.split("\n")
        first_run = frame.paragraphs[0].runs[0]
        first_run.text = lines[0]

        for extra in frame.paragraphs[0].runs[1:]:
            extra._element.getparent().remove(extra._element)

        for line in lines[1:]:
            p = frame.add_paragraph()
            run = p.add_run()
            run.text = line
            _copy_run_formatting(first_run.font, run.font)

        # ── Overflow protection ──────────────────────────────────────
        # If the text overflows the shape's height, iteratively shrink
        # font sizes until the content fits (floor: 8 pt).
        _auto_fit_text_frame(frame, shape, value)


# ─── Overflow-protection helper ───────────────────────────────────────────────

# Minimum font size (in pt) before we stop shrinking.  Going below this
# produces unreadable text; at that point it is better to let the content
# overflow slightly than to render illegible slides.
_MIN_FONT_PT = 8.0


def _auto_fit_text_frame(frame, shape, text_value: str) -> None:
    """Iteratively shrink font sizes in *frame* until the text fits *shape*.

    Reuses :func:`app.engines.canvas_calculator.font_metrics.estimate_box`
    for height estimation so the template-injector and canvas-fallback
    paths stay consistent.

    Args:
        frame: The ``TextFrame`` that was just populated by ``_write_text``.
        shape: The parent ``BaseShape`` (provides ``shape.height``).
        text_value: The original text written into the frame.
    """
    from pptx.util import Emu, Pt

    from app.engines.canvas_calculator.font_metrics import estimate_box

    available_h = shape.height  # EMU
    if available_h is None or available_h <= 0:
        return  # no height constraint — nothing to do

    # Account for the text-frame's internal vertical margins so the
    # usable area is accurate.
    try:
        top_m = frame.margin_top or 0
        btm_m = frame.margin_bottom or 0
        available_h -= (top_m + btm_m)
    except Exception:
        pass

    if available_h <= 0:
        return

    # Also account for horizontal margins when estimating chars-per-line.
    try:
        left_m = frame.margin_left or 0
        right_m = frame.margin_right or 0
        available_w = shape.width - (left_m + right_m)
    except Exception:
        available_w = shape.width

    if available_w is None or available_w <= 0:
        available_w = shape.width or Emu(0)

    # Collect every run so we can shrink them all.
    runs = []
    for paragraph in frame.paragraphs:
        for run in paragraph.runs:
            runs.append(run)

    if not runs:
        return

    # Determine starting font size from the first run.
    first_size = runs[0].font.size
    if first_size is None:
        return  # no font size set — cannot auto-fit

    current_pt = first_size / 12700  # EMU → pt
    if current_pt <= _MIN_FONT_PT:
        return  # already at floor

    # Shrink loop: reduce font size by 0.5 pt each iteration until
    # estimated text height ≤ available_h or we hit the floor.
    # Maximum iterations = (current_pt - _MIN_FONT_PT) / 0.5, capped at 60.
    max_iterations = min(60, int((current_pt - _MIN_FONT_PT) / 0.5) + 1)

    for _ in range(max_iterations):
        if current_pt <= _MIN_FONT_PT:
            break

        estimated_h = estimate_box(text_value, current_pt, available_w)
        if estimated_h <= available_h:
            break  # fits!
        current_pt -= 0.5

    # Apply the final font size to every run.
    final_size = Pt(current_pt)
    for run in runs:
        run.font.size = final_size

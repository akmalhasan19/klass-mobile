"""Placeholder resolver for the Master Template Injection pipeline.

Given a target ``python-pptx`` slide and a ``shape_name`` (as declared in the
template manifest), the resolver locates the corresponding ``BaseShape`` so
the injector can write content into it.

Key behaviours
--------------
- Resolution is by ``shape.name`` (the same value the manifest stores under
  ``PlaceholderSpec.shape_name``).
- Shapes nested inside **group shapes** are searched recursively — a bare
  ``for shape in slide.shapes`` would miss them, which is a common failure
  mode with designer-built masters.
- A missing placeholder resolves to ``None`` (never raises) and is reported
  both via ``logging`` and, optionally, via a caller-supplied ``warnings``
  list so the injector can surface it in its ``InjectionResult``.

Kind handling
-------------
``PlaceholderKind`` is ``"text"`` (supported today) or ``"image"`` (resolved
by name, but its population is deferred to a future injector revision — image
fill needs asset handling that doesn't exist yet).  Resolution itself is
kind-agnostic; the distinction only matters when *filling* the shape, which
remains the injector's responsibility.
"""
from __future__ import annotations

import logging
from typing import Iterable

# pyrefly: ignore [missing-import]
from pptx.enum.shapes import MSO_SHAPE_TYPE
# pyrefly: ignore [missing-import]
from pptx.shapes.base import BaseShape
# pyrefly: ignore [missing-import]
from pptx.slide import Slide

from app.engines.pptx_injector.manifest import PlaceholderSpec

logger = logging.getLogger("klass-media-generator")

# Kinds the injection pipeline can currently populate.  "image" is resolved
# by name but not yet filled, hence it is intentionally excluded here.
SUPPORTED_FILL_KINDS: frozenset[str] = frozenset({"text"})


def _iter_shapes(shapes: Iterable[BaseShape]) -> Iterable[BaseShape]:
    """Yield *shapes* and, recursively, any shapes nested inside groups."""
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            nested = getattr(shape, "shapes", None)
            if nested is not None:
                yield from _iter_shapes(nested)


def resolve_shape(
    slide: Slide,
    shape_name: str,
    *,
    warnings: list[str] | None = None,
) -> BaseShape | None:
    """Resolve a master shape on *slide* by its ``shape.name``.

    Searches top-level shapes and shapes nested inside group shapes.

    Args:
        slide: The slide (or duplicated master slide) to search.
        shape_name: The ``shape.name`` to match, taken from the manifest's
            ``PlaceholderSpec.shape_name``.
        warnings: Optional list to which a human-readable warning is appended
            when the shape cannot be found.  Always logged regardless.

    Returns:
        The first matching :class:`~pptx.shapes.base.BaseShape`, or ``None``
        if no shape with that name exists.
    """
    for shape in _iter_shapes(slide.shapes):
        if shape.name == shape_name:
            return shape

    message = f"Placeholder shape '{shape_name}' not found on slide."
    logger.warning(message)
    if warnings is not None:
        warnings.append(message)
    return None


def resolve_placeholder(
    slide: Slide,
    spec: PlaceholderSpec,
    *,
    warnings: list[str] | None = None,
) -> BaseShape | None:
    """Resolve a placeholder described by a :class:`PlaceholderSpec`.

    Validates *spec.kind* (only ``"text"`` is fillable today; ``"image"`` is
    accepted for resolution but its population is deferred) and then resolves
    the underlying shape via :func:`resolve_shape`.

    Raises:
        ValueError: If *spec.kind* is neither supported nor the reserved
            ``"image"`` kind.
    """
    if spec.kind not in SUPPORTED_FILL_KINDS and spec.kind != "image":
        raise ValueError(f"Unsupported placeholder kind: {spec.kind!r}")

    return resolve_shape(slide, spec.shape_name, warnings=warnings)

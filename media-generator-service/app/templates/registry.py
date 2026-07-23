"""Template registry for all master templates.

Loads master templates at startup and provides a unified lookup API keyed by
``template_id`` (which equals ``SlideBlueprint.theme_id``):

- **PPTX** — master ``.pptx`` + manifest JSON (Template Injector pipeline).
- **HTML** — self-contained Jinja2 master ``.html`` (PDF + WebView preview
  pipeline, Fase 2).
- **DOCX** — master ``.docx`` with ``docxtpl`` placeholders (DOCX pipeline,
  Fase 1).

All three masters share the same ``template_id`` so a single registry entry
describes every format for a given visual design.  This keeps one source of
truth instead of a parallel registry per format.

Design decisions:
- Manifests are immutable JSON cached in memory (share across requests).
- Master files are opened per request (``python-pptx`` / ``docxtpl`` are not
  thread-safe), so we only store the file paths, not the loaded documents.
- Validation happens at startup — fail fast with
  ``ServiceMisconfiguredError`` if a required master is missing or (for PPTX)
  a manifest references a shape that doesn't exist in the master template.

Master contract (read this before writing the injector)
-------------------------------------------------------
The ``slide_index`` field in a manifest ``LayoutManifest`` indexes the
master ``.pptx`` ``Presentation.slides`` **collection**, *not*
``slide_layouts``.  This master ships three pre-designed slides
(``title`` / ``content`` / ``assessment``) whose placeholder shapes
(``"Title 1"``, ``"Objectives Box 3"``, ``"Content Placeholder 3"``, …)
live on those slides — not on the bare slide layouts.

Consequently the injector must **duplicate the source slide**
(``source = Presentation(master_path).slides[layout.slide_index]``) into
the output deck and then fill placeholders by ``shape.name``.  Using
``presentation.slides.add_slide(layout.slide_layout)`` (as the plan's
prototype suggested) would NOT yield these designed shapes and must not
be used with this template.
"""
from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path

# pyrefly: ignore [missing-import]
from pptx import Presentation

from app.engines.pptx_injector.manifest import TemplateManifest, load_manifest
from app.errors import ServiceMisconfiguredError

logger = logging.getLogger("klass-media-generator")


@dataclass(frozen=True)
class TemplateEntry:
    """A loaded template: manifest + paths to all format masters.

    This is the registry's public payload.  Every format master for a given
    ``template_id`` is referenced here:

    >>> entry = registry.get("klass-educational-v1")
    >>> entry.master_path      # Path to the .pptx (open per request)
    >>> entry.manifest         # Cached, immutable TemplateManifest
    >>> entry.html_master_path # Path to the Jinja2 .html master
    >>> entry.docx_master_path # Path to the .docx master (docxtpl)

    Attributes:
        manifest: Parsed and validated manifest (cached, shared across requests).
        master_path: Absolute path to the master ``.pptx`` file.
        html_master_path: Absolute path to the Jinja2 HTML master (PDF/preview).
        docx_master_path: Absolute path to the ``.docx`` master (DOCX pipeline).
    """

    manifest: TemplateManifest
    master_path: Path
    html_master_path: Path | None = None
    docx_master_path: Path | None = None


@dataclass
class TemplateRegistry:
    """Registry of all available master templates.

    Populated at startup via :meth:`load_templates`.  The registry is
    intended to be a singleton stored in the FastAPI lifespan and
    injected into the PPTX generator.

    Usage::

        registry = TemplateRegistry()
        registry.load_templates/templates_dir)

        entry = registry.get("klass-educational-v1")
        # entry.master_path → Path to .pptx
        # entry.manifest → TemplateManifest
    """

    _templates: dict[str, TemplateEntry] = field(default_factory=dict)

    def load_templates(self, templates_dir: Path) -> None:
        """Discover and load all templates from *templates_dir*.

        Expected directory structure::

            templates_dir/
                masters/
                    klass-educational-v1.pptx
                manifests/
                    klass-educational-v1.json

        Raises:
            ServiceMisconfiguredError: If no templates are found, or if
                a manifest references shape names that don't exist in
                the master template.
        """
        masters_dir = templates_dir / "masters"
        manifests_dir = templates_dir / "manifests"

        if not masters_dir.is_dir() or not manifests_dir.is_dir():
            raise ServiceMisconfiguredError(
                "Template directories not found.",
                {
                    "masters_dir": str(masters_dir),
                    "manifests_dir": str(manifests_dir),
                },
            )

        # Discover manifests by *.json files
        manifest_files = sorted(manifests_dir.glob("*.json"))
        if not manifest_files:
            raise ServiceMisconfiguredError(
                "No manifest files found in templates/manifests/.",
                {"manifests_dir": str(manifests_dir)},
            )

        for manifest_file in manifest_files:
            template_id = manifest_file.stem
            master_file = masters_dir / f"{template_id}.pptx"

            if not master_file.is_file():
                raise ServiceMisconfiguredError(
                    f"Master template not found for manifest '{template_id}'.",
                    {
                        "template_id": template_id,
                        "expected_master": str(master_file),
                    },
                )

            # Load and validate manifest
            manifest = load_manifest(manifest_file)

            # The manifest's own template_id must agree with the filename
            # stem — the registry keys entries by stem, so a mismatch would
            # make the manifest unreachable under its declared id.
            if manifest.template_id != template_id:
                raise ServiceMisconfiguredError(
                    "Manifest template_id does not match its filename.",
                    {
                        "expected_template_id": template_id,
                        "manifest_template_id": manifest.template_id,
                        "manifest_file": str(manifest_file),
                    },
                )

            # Validate shape names against master
            self._validate_shape_names(template_id, master_file, manifest)

            # Discover the HTML (PDF/preview) and DOCX masters for the same
            # template_id.  They are required for the migration's gate
            # ("all master templates ready"), so missing files fail fast.
            html_master = masters_dir / f"{template_id}.html"
            docx_master = masters_dir / f"{template_id}.docx"

            if not html_master.is_file():
                raise ServiceMisconfiguredError(
                    f"HTML master template not found for template '{template_id}'.",
                    {
                        "template_id": template_id,
                        "expected_master": str(html_master),
                    },
                )
            if not docx_master.is_file():
                raise ServiceMisconfiguredError(
                    f"DOCX master template not found for template '{template_id}'.",
                    {
                        "template_id": template_id,
                        "expected_master": str(docx_master),
                    },
                )

            self._templates[template_id] = TemplateEntry(
                manifest=manifest,
                master_path=master_file,
                html_master_path=html_master,
                docx_master_path=docx_master,
            )
            logger.info(
                "Loaded template '%s' (v%s) from %s",
                template_id,
                manifest.version,
                master_file,
            )

    def get(self, template_id: str) -> TemplateEntry:
        """Retrieve a template by its ID.

        Returns:
            TemplateEntry: a frozen payload holding ``master_path`` (the
            absolute path to the ``.pptx`` — reopen per request) and
            ``manifest`` (the cached, immutable ``TemplateManifest``).
            This is the equivalent of the plan's ``(master_path, manifest)``
            tuple, exposed as a typed dataclass.

        Raises:
            KeyError: If no template with the given ID is registered.
        """
        entry = self._templates.get(template_id)
        if entry is None:
            available = list(self._templates.keys())
            raise KeyError(
                f"Template '{template_id}' not found. "
                f"Available templates: {available}"
            )
        return entry

    @property
    def template_ids(self) -> list[str]:
        """Return all registered template IDs."""
        return list(self._templates.keys())

    def get_html_master(self, template_id: str) -> Path:
        """Return the absolute path to the Jinja2 HTML master for *template_id*.

        Raises:
            KeyError: If no template with the given ID is registered.
            ServiceMisconfiguredError: If the HTML master is missing (should not
                happen — ``load_templates`` validates presence at startup).
        """
        entry = self.get(template_id)
        if entry.html_master_path is None or not entry.html_master_path.is_file():
            raise ServiceMisconfiguredError(
                f"HTML master template missing for template '{template_id}'.",
                {"template_id": template_id},
            )
        return entry.html_master_path

    def get_docx_master(self, template_id: str) -> Path:
        """Return the absolute path to the ``.docx`` master for *template_id*.

        Raises:
            KeyError: If no template with the given ID is registered.
            ServiceMisconfiguredError: If the DOCX master is missing (should not
                happen — ``load_templates`` validates presence at startup).
        """
        entry = self.get(template_id)
        if entry.docx_master_path is None or not entry.docx_master_path.is_file():
            raise ServiceMisconfiguredError(
                f"DOCX master template missing for template '{template_id}'.",
                {"template_id": template_id},
            )
        return entry.docx_master_path

    def _validate_shape_names(
        self,
        template_id: str,
        master_path: Path,
        manifest: TemplateManifest,
    ) -> None:
        """Validate that all shape names in the manifest exist in the master.

        Raises:
            ServiceMisconfiguredError: If any shape name is missing.
        """
        prs = Presentation(str(master_path))
        missing: list[dict[str, str]] = []

        for layout in manifest.slide_layouts:
            if layout.slide_index >= len(prs.slides):
                missing.append({
                    "layout_id": layout.layout_id,
                    "slide_index": str(layout.slide_index),
                    "error": f"slide_index {layout.slide_index} out of range "
                             f"(master has {len(prs.slides)} slides)",
                })
                continue

            slide = prs.slides[layout.slide_index]
            shape_names = {shape.name for shape in slide.shapes}

            for placeholder in layout.placeholders:
                if placeholder.shape_name not in shape_names:
                    missing.append({
                        "layout_id": layout.layout_id,
                        "placeholder_id": placeholder.placeholder_id,
                        "shape_name": placeholder.shape_name,
                        "available_shapes": str(sorted(shape_names)),
                    })

        if missing:
            raise ServiceMisconfiguredError(
                f"Shape name mismatch in template '{template_id}'. "
                "Master template and manifest are out of sync.",
                {
                    "template_id": template_id,
                    "master_path": str(master_path),
                    "mismatches": missing,
                },
            )

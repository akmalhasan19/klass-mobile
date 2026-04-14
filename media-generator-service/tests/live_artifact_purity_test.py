"""
Live Artifact Regeneration Test
================================
This script generates DOCX, PDF, and PPTX artifacts using the actual
generators and inspects each for forbidden metadata/scaffold content.

Expected Outcome:
  - ✅ Ada: summary paragraph (tanpa label), learning objectives, section titles, activities
  - ❌ Tidak ada: context table, asset list, teacher notes, section purpose subtitles
  - ❌ Tidak ada: scaffold prose ("Bagian ini disusun untuk...", etc.)
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from docx import Document
from pptx import Presentation

# -- Bootstrap -----------------------------------------------------------------
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from app.document_model import build_render_document
from app.generators.docx_generator import DocxGenerator
from app.generators.pdf_generator import PdfGenerator
from app.generators.pptx_generator import PptxGenerator
from app.settings import get_settings
from tests.helpers import cleanup_artifact, sample_request

# -- Forbidden patterns --------------------------------------------------------
FORBIDDEN_LABELS = [
    "Ringkasan",          # Summary label
    "Konteks Kelas",      # Context table
    "Catatan untuk Guru",  # Teacher notes
    "Aset Pendukung",     # Asset list
    "Classroom Context",
    "Teacher Notes",
    "Suggested Assets",
]

SCAFFOLD_PATTERNS = [
    r"bagian ini disusun untuk",
    r"fokus utamanya meliputi",
    r"jelaskan ide pokoknya secara runtut",
    r"sampaikan inti materinya secara singkat",
    r"dorong siswa merangkum kembali",
    r"this section is written for",
    r"the main focus includes",
    r"present the main idea in sequence",
    r"keep the explanation concise, clear, and ready for presentation",
    r"encourage students to restate the key idea",
]

INTERNAL_PATTERNS = [
    r"schema_version",
    r"Return exactly one JSON object",
    r"body_blocks",
    r"media_content_draft\.v1",
    r"media_prompt_understanding\.v1",
]

# Section purpose texts from the sample request that should NOT appear
SECTION_PURPOSE_TEXTS = [
    "Menjelaskan target belajar utama sebelum latihan dimulai.",
    "Memberi contoh singkat lalu latihan mandiri.",
]

# -- Expected content ----------------------------------------------------------
EXPECTED_CONTENT = [
    "Handout Pecahan Kelas 5",                                              # Title
    "Handout ringkas untuk memperkenalkan pecahan senilai dan latihan dasar.",  # Summary
    "Tujuan Belajar",                                                       # Section title
    "Contoh dan Latihan",                                                   # Section title
    "Latihan Mandiri",                                                      # Activity title
]


def check_text(text: str, label: str) -> list[str]:
    """Check a text blob for forbidden content. Return list of failure messages."""
    failures: list[str] = []

    for forbidden in FORBIDDEN_LABELS:
        if forbidden in text:
            failures.append(f"[{label}] Found forbidden label: '{forbidden}'")

    for pattern in SCAFFOLD_PATTERNS:
        if re.search(pattern, text, re.IGNORECASE):
            failures.append(f"[{label}] Found scaffold pattern: '{pattern}'")

    for pattern in INTERNAL_PATTERNS:
        if re.search(pattern, text, re.IGNORECASE):
            failures.append(f"[{label}] Found internal pattern: '{pattern}'")

    for purpose in SECTION_PURPOSE_TEXTS:
        if purpose in text:
            failures.append(f"[{label}] Found section purpose text leaked into body: '{purpose}'")

    return failures


def test_docx_artifact() -> tuple[list[str], list[str]]:
    """Generate and inspect DOCX artifact."""
    request = sample_request("docx")
    render_doc = build_render_document(request.generation_spec)
    generator = DocxGenerator()
    metadata = generator.generate(request, render_doc, get_settings())

    document = Document(metadata["artifact_locator"]["value"])
    full_text = "\n".join(p.text for p in document.paragraphs if p.text.strip())

    failures = check_text(full_text, "DOCX")
    found = []

    # Check expected content is present
    for expected in EXPECTED_CONTENT:
        if expected in full_text:
            found.append(f"  ✅ Found expected: '{expected}'")
        else:
            failures.append(f"[DOCX] Missing expected content: '{expected}'")

    # Check no table elements (context table)
    if document.tables:
        failures.append(f"[DOCX] Found {len(document.tables)} table(s) - context table should be removed")
    else:
        found.append("  ✅ No tables found (context table removed)")

    cleanup_artifact(metadata)
    return failures, found


def test_pdf_artifact() -> tuple[list[str], list[str]]:
    """Generate and inspect PDF artifact."""
    request = sample_request("pdf")
    render_doc = build_render_document(request.generation_spec)
    generator = PdfGenerator()
    metadata = generator.generate(request, render_doc, get_settings())

    artifact_path = Path(metadata["artifact_locator"]["value"])
    raw_bytes = artifact_path.read_bytes()

    failures = []
    found = []

    # PDF starts with %PDF
    if raw_bytes.startswith(b"%PDF"):
        found.append("  ✅ Valid PDF header")
    else:
        failures.append("[PDF] Invalid PDF header")

    # Page count
    page_count = metadata.get("page_count", 0)
    if page_count and page_count >= 1:
        found.append(f"  ✅ Page count: {page_count}")
    else:
        failures.append(f"[PDF] Unexpected page count: {page_count}")

    # Check raw bytes for forbidden patterns
    for pattern in INTERNAL_PATTERNS:
        if re.search(pattern.encode(), raw_bytes, re.IGNORECASE):
            failures.append(f"[PDF] Found internal pattern in raw bytes: '{pattern}'")

    # Check forbidden labels in raw bytes
    for forbidden in FORBIDDEN_LABELS:
        if forbidden.encode("utf-8") in raw_bytes:
            failures.append(f"[PDF] Found forbidden label in raw bytes: '{forbidden}'")

    if not any("[PDF]" in f for f in failures):
        found.append("  ✅ No forbidden labels/patterns in PDF bytes")

    # Check expected content in raw bytes
    for expected in ["Handout Pecahan Kelas 5", "Tujuan"]:
        if expected.encode("utf-8") in raw_bytes:
            found.append(f"  ✅ Found expected in PDF: '{expected}'")

    cleanup_artifact(metadata)
    return failures, found


def test_pptx_artifact() -> tuple[list[str], list[str]]:
    """Generate and inspect PPTX artifact."""
    request = sample_request("pptx")
    render_doc = build_render_document(request.generation_spec)
    generator = PptxGenerator()
    metadata = generator.generate(request, render_doc, get_settings())

    presentation = Presentation(metadata["artifact_locator"]["value"])
    all_text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = getattr(shape, "text", "")
                if isinstance(text, str) and text.strip():
                    all_text.append(text.strip())
    full_text = "\n".join(all_text)

    failures = check_text(full_text, "PPTX")
    found = []

    # Check expected content
    pptx_expected = [
        "Handout Pecahan Kelas 5",
        "Tujuan Belajar",                   # Learning objectives heading
        "Aktivitas dan Penilaian",          # Activity slide
    ]
    for expected in pptx_expected:
        if expected in full_text:
            found.append(f"  ✅ Found expected: '{expected}'")
        else:
            failures.append(f"[PPTX] Missing expected content: '{expected}'")

    # Slide count
    slide_count = metadata.get("slide_count", 0)
    if slide_count and slide_count >= 3:
        found.append(f"  ✅ Slide count: {slide_count}")
    else:
        failures.append(f"[PPTX] Unexpected slide count: {slide_count}")

    # Verify learning objectives panel exists (not context/metadata panel)
    assert any("Tujuan" in t for t in all_text), "Learning objectives panel missing from title slide"
    found.append("  ✅ Learning objectives panel on title slide (not metadata context)")

    cleanup_artifact(metadata)
    return failures, found


def main():
    print("=" * 70)
    print("  LIVE ARTIFACT REGENERATION TEST")
    print("  Verifying artifact purity after metadata/scaffold removal")
    print("=" * 70)

    all_failures: list[str] = []
    formats = [
        ("DOCX", test_docx_artifact),
        ("PDF", test_pdf_artifact),
        ("PPTX", test_pptx_artifact),
    ]

    for format_name, test_fn in formats:
        print(f"\n{'─' * 50}")
        print(f"  Testing {format_name} artifact...")
        print(f"{'─' * 50}")
        try:
            failures, found = test_fn()
            for f in found:
                print(f)
            if failures:
                for f in failures:
                    print(f"  ❌ {f}")
                all_failures.extend(failures)
            else:
                print(f"\n  🎉 {format_name}: ALL CHECKS PASSED")
        except Exception as e:
            msg = f"[{format_name}] Exception during test: {e}"
            print(f"  ❌ {msg}")
            all_failures.append(msg)

    print(f"\n{'=' * 70}")
    if all_failures:
        print(f"  ❌ FAILED: {len(all_failures)} issue(s) found")
        for f in all_failures:
            print(f"    • {f}")
        print(f"{'=' * 70}")
        sys.exit(1)
    else:
        print("  ✅ ALL FORMATS PASSED - Artifacts are clean!")
        print("  No metadata, scaffold, context tables, asset lists,")
        print("  teacher notes, or section purposes found in final artifacts.")
        print(f"{'=' * 70}")
        sys.exit(0)


if __name__ == "__main__":
    main()

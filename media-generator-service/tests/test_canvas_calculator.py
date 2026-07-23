"""Unit tests for the Canvas Calculator fallback engine (Fase 3).

Tests the geometry computation (``CanvasLayoutEngine``), font metrics
(``estimate_box``), shape rendering (``CanvasShapeRenderer``), and
the overall ``render_slide`` pipeline end-to-end.

Coverage
--------
* ``_pick_columns`` heuristic — correct column count for 1–10+ cards.
* ``compute_layout`` — correct box geometry (count, position, size).
* ``render_slide`` — shapes are actually added to a ``Presentation``.
* ``estimate_box`` — returns positive, reasonable EMU heights.
* Shape renderer — ``draw_title`` and ``draw_card`` add visible content.
* Overflow guard — ``compute_layout`` raises when cards exceed slide height.
* CanvasShapeRenderer — uses correct colours, fonts, and shape types.
"""
from __future__ import annotations

import math
from pathlib import Path

import pytest
# pyrefly: ignore [missing-import]
from pptx import Presentation
# pyrefly: ignore [missing-import]
from pptx.util import Emu, Inches, Pt

from app.engines.blueprint import Card, ContentBlock, Slide
from app.engines.canvas_calculator.font_metrics import estimate_box
from app.engines.canvas_calculator.layout_engine import (
    Box,
    CanvasLayoutEngine,
    SlideLayout,
)
from app.engines.canvas_calculator.shape_renderer import (
    CanvasShapeRenderer,
    _font_sizes,
)
from app.errors import GenerationError

# ---------------------------------------------------------------------------
# Fixtures / shared builders
# ---------------------------------------------------------------------------


def _card(text: str, kind: str = "paragraph") -> Card:
    return Card(body_blocks=[ContentBlock(kind=kind, content=text)])


def _content_slide(title: str, card_count: int) -> Slide:
    return Slide(
        slide_type="content",
        title=title,
        cards=[_card(f"Card {i}") for i in range(card_count)],
    )


# ---------------------------------------------------------------------------
# 1. Column heuristic
# ---------------------------------------------------------------------------


def test_pick_columns_one_card() -> None:
    engine = CanvasLayoutEngine()
    assert engine._pick_columns(1) == 1


def test_pick_columns_two_cards() -> None:
    engine = CanvasLayoutEngine()
    assert engine._pick_columns(2) == 2


def test_pick_columns_three_four_cards() -> None:
    engine = CanvasLayoutEngine()
    assert engine._pick_columns(3) == 2
    assert engine._pick_columns(4) == 2


def test_pick_columns_five_to_nine_cards() -> None:
    engine = CanvasLayoutEngine()
    for n in range(5, 10):
        assert engine._pick_columns(n) == 3, f"expected 3 cols for {n} cards"


def test_pick_columns_ten_or_more_cards() -> None:
    engine = CanvasLayoutEngine()
    for n in range(10, 15):
        assert engine._pick_columns(n) == 4, f"expected 4 cols for {n} cards"


# ---------------------------------------------------------------------------
# 2. Layout geometry
# ---------------------------------------------------------------------------


def test_compute_layout_returns_correct_number_of_boxes() -> None:
    engine = CanvasLayoutEngine()
    slide = _content_slide("Geometry", 6)
    layout = engine.compute_layout(slide)

    assert isinstance(layout, SlideLayout)
    assert len(layout.card_boxes) == 6, f"expected 6 boxes, got {len(layout.card_boxes)}"


def test_compute_layout_title_box_at_top() -> None:
    engine = CanvasLayoutEngine()
    slide = _content_slide("Title", 3)
    layout = engine.compute_layout(slide)

    title_box = layout.title_box
    assert title_box.y == engine._margin
    assert title_box.h == engine._title_band
    assert title_box.x == engine._margin


def test_compute_layout_boxes_do_not_overlap() -> None:
    engine = CanvasLayoutEngine()
    slide = _content_slide("No Overlap", 6)
    layout = engine.compute_layout(slide)

    for i, a in enumerate(layout.card_boxes):
        for j, b in enumerate(layout.card_boxes):
            if i >= j:
                continue
            # Check no horizontal overlap
            a_right = a.x + a.w
            b_right = b.x + b.w
            if a_right > b.x and b_right > a.x:
                # Same row — check vertical non-overlap or same-y
                a_bottom = a.y + a.h
                b_bottom = b.y + b.h
                overlap = a_bottom > b.y and b_bottom > a.y
                assert not overlap, f"boxes {i} and {j} overlap: {a} vs {b}"


def test_compute_layout_with_single_card() -> None:
    engine = CanvasLayoutEngine()
    slide = _content_slide("Solo", 1)
    layout = engine.compute_layout(slide)

    assert len(layout.card_boxes) == 1
    box = layout.card_boxes[0]
    assert box.w > Emu(0)
    assert box.h > Emu(0)


# ---------------------------------------------------------------------------
# 3. render_slide (end-to-end with real ShapeRenderer)
# ---------------------------------------------------------------------------


def test_render_slide_adds_shapes_to_presentation() -> None:
    engine = CanvasLayoutEngine(renderer=CanvasShapeRenderer())
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = _content_slide("Rendered", 4)
    layout = engine.render_slide(prs, slide)

    assert len(prs.slides) == 1
    rendered_slide = prs.slides[0]
    shape_count = len(rendered_slide.shapes)

    # At minimum: 1 title box + 4 card shapes = 5 shapes
    assert shape_count >= 5, f"expected ≥5 shapes, got {shape_count}"


def test_render_slide_adds_cards_as_shapes() -> None:
    """The rendered slide has at least one shape per card plus a title."""
    engine = CanvasLayoutEngine(renderer=CanvasShapeRenderer())
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = _content_slide("Cards", 3)
    engine.render_slide(prs, slide)

    rendered_slide = prs.slides[0]
    shape_count = len(rendered_slide.shapes)
    # 1 title box + 3 cards = at least 4 shapes (blank layout may have 0)
    assert shape_count >= 4, f"expected ≥4 shapes for title+3 cards, got {shape_count}"


def test_render_slide_title_has_text() -> None:
    engine = CanvasLayoutEngine(renderer=CanvasShapeRenderer())
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = _content_slide("Hello World", 1)
    engine.render_slide(prs, slide)

    rendered_slide = prs.slides[0]
    all_text = " ".join(
        s.text for s in rendered_slide.shapes if hasattr(s, "text") and s.text
    )
    assert "Hello World" in all_text, f"title not found in shapes: {all_text[:200]}"


def test_render_slide_with_heading_in_card() -> None:
    engine = CanvasLayoutEngine(renderer=CanvasShapeRenderer())
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = Slide(
        slide_type="content",
        title="Slide Judul",
        cards=[
            Card(
                heading="Card Heading",
                body_blocks=[ContentBlock(kind="paragraph", content="Body text here.")],
            ),
        ],
    )
    engine.render_slide(prs, slide)

    rendered_slide = prs.slides[0]
    card_texts = " ".join(
        s.text for s in rendered_slide.shapes if hasattr(s, "text") and s.text
    )
    assert "Card Heading" in card_texts, f"heading not found: {card_texts[:200]}"
    assert "Body text here." in card_texts


# ---------------------------------------------------------------------------
# 4. Font metrics
# ---------------------------------------------------------------------------


def test_estimate_box_returns_positive_emu() -> None:
    height = estimate_box("Hello world", 12, Inches(4))
    assert height > Emu(0)


def test_estimate_box_scales_with_text_length() -> None:
    short_h = estimate_box("short", 12, Inches(4))
    long_h = estimate_box("longer " * 50, 12, Inches(4))
    assert long_h > short_h, "longer text should need more height"


def test_estimate_box_scales_with_font_size() -> None:
    small_h = estimate_box("text", 10, Inches(4))
    large_h = estimate_box("text", 24, Inches(4))
    assert large_h > small_h, "larger font should need more height"


def test_estimate_box_newline_increases_height() -> None:
    single_h = estimate_box("one line", 12, Inches(4))
    multi_h = estimate_box("line1\nline2", 12, Inches(4))
    assert multi_h > single_h, "explicit line breaks should increase height"


# ---------------------------------------------------------------------------
# 5. Shape renderer colours and types
# ---------------------------------------------------------------------------


def test_shape_renderer_draw_title_creates_textbox() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    renderer = CanvasShapeRenderer()
    box = Box(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    renderer.draw_title(slide, "Title Here", box)

    # Should have at least one shape with text
    shape_texts = [s.text for s in slide.shapes if hasattr(s, "text")]
    assert any("Title Here" in t for t in shape_texts), f"title text missing: {shape_texts}"


def test_shape_renderer_draw_card_adds_shape_with_text() -> None:
    """``draw_card`` adds at least one shape containing the card's text."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    renderer = CanvasShapeRenderer()
    box = Box(Inches(0.5), Inches(0.5), Inches(4), Inches(3))
    card = Card(body_blocks=[ContentBlock(kind="paragraph", content="Card text")])
    renderer.draw_card(slide, card, box)

    # At minimum the card text should appear in a shape.
    shape_texts = [s.text for s in slide.shapes if hasattr(s, "text")]
    assert any("Card text" in t for t in shape_texts), (
        f"card text not found in shapes: {shape_texts}"
    )


def test_font_sizes_scaling() -> None:
    # When box is very wide → large font
    h, b = _font_sizes(Inches(7))
    assert h >= 16 and b >= 12, f"wide box should have large font: ({h}, {b})"

    # When box is narrow → small font
    h, b = _font_sizes(Inches(2))
    assert h <= 14 and b <= 12, f"narrow box should have small font: ({h}, {b})"


# ---------------------------------------------------------------------------
# 6. Overflow guard
# ---------------------------------------------------------------------------


def test_compute_layout_guards_against_extreme_overflow() -> None:
    """Slides whose layout exceeds slide height raise ``GenerationError``."""
    engine = CanvasLayoutEngine()
    # 50 cards each with long text — guarantees overflow of 7.5" slide.
    slide = Slide(
        slide_type="content",
        title="Overflow",
        cards=[_card("X" * 200) for _ in range(50)],
    )

    with pytest.raises(GenerationError):
        engine.compute_layout(slide)


# ---------------------------------------------------------------------------
# 7. File integrity
# ---------------------------------------------------------------------------


def test_canvas_output_is_valid_pptx() -> None:
    """The canvas output is a valid PPTX that re-opens cleanly."""
    import tempfile
    import os

    engine = CanvasLayoutEngine(renderer=CanvasShapeRenderer())
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = _content_slide("Valid File", 2)
    engine.render_slide(prs, slide)

    fd, path = tempfile.mkstemp(suffix=".pptx", prefix="test_canvas_")
    os.close(fd)
    try:
        prs.save(path)
        # Re-open
        reopened = Presentation(path)
        assert len(reopened.slides) == 1
    finally:
        Path(path).unlink(missing_ok=True)

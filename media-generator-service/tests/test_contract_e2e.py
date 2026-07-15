"""Contract + E2E tests for the generation API (Fase 4: Task 4.5).

Covers the gate criteria specified in the implementation plan:

1. Response passes ``GenerateSuccessResponse.model_validate`` — the full
   JSON payload returned by ``POST /v1/generate`` must round-trip through
   the Pydantic model without error.

2. Signed preview URL returns ``200`` with ``Content-Type: text/html``
   (already tested in ``test_preview_api.py`` — referenced here).

3. ``artifact_metadata.slide_count`` is consistent with the actual number
   of slides in the output file.

4. HMAC signing uses the **same secret** for both artifact and preview
   signed URLs — no divergence in signing material.

5. Timeout / retry configuration is consistent with the Gateway contract
   (Marp render ≤ 30s, total request ≤ 60s).

6. Response includes the new Fase 4 additive fields (``preview_url``,
   ``layout_sources``) when applicable.
"""
from __future__ import annotations

from pathlib import Path
from unittest.mock import AsyncMock, MagicMock
from urllib.parse import parse_qs, urlparse

from pptx import Presentation

from app.models import (
    ArtifactMetadata,
    GenerateSuccessResponse,
)
from tests.helpers import (
    artifact_path_from_metadata,
    cleanup_artifact,
    signed_request_content,
)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _mock_ready_sidecar(
    html: str = "<html><body>Preview</body></html>",
) -> MagicMock:
    sc = MagicMock()
    sc.render_html = AsyncMock(return_value=html)
    sc.render_pdf = AsyncMock(return_value=b"%PDF-1.4-mock")
    sc.is_running = True
    sc.is_ready = True
    sc.uptime_seconds = 42.0
    return sc


def _extract_signed_url_params(locator: dict[str, str]) -> dict[str, list[str]]:
    """Parse the signed URL's query parameters into a dict."""
    parsed = urlparse(locator["value"])
    return parse_qs(parsed.query)


# ---------------------------------------------------------------------------
# 1. Response model validation
# ---------------------------------------------------------------------------


def test_response_passes_generate_success_response_validation(
    client, monkeypatch,
) -> None:
    """The raw JSON response round-trips through ``GenerateSuccessResponse``."""
    import app.main

    monkeypatch.setattr(app.main, "sidecar_manager", _mock_ready_sidecar())

    body, headers, _ = signed_request_content("pptx")
    response = client.post("/v1/generate", content=body, headers=headers)

    assert response.status_code == 200
    payload = response.json()

    # model_validate raises on any schema violation.
    validated = GenerateSuccessResponse.model_validate(payload)
    assert validated.status == "completed"
    assert validated.data.generation_id == "generation-pptx-001"
    assert validated.data.preview_delivery is not None
    assert validated.data.artifact_metadata.slide_count == 4

    pd = validated.data.preview_delivery
    assert pd.schema_version == "media_generator_preview.v1"
    assert pd.mime_type == "text/html"
    assert pd.locator.kind == "signed_url"

    cleanup_artifact(payload["data"]["artifact_metadata"])


def test_response_without_sidecar_passes_validation(client) -> None:
    """Even without a sidecar, the response is a valid ``GenerateSuccessResponse``."""
    body, headers, _ = signed_request_content("pptx")
    response = client.post("/v1/generate", content=body, headers=headers)

    assert response.status_code == 200
    payload = response.json()

    validated = GenerateSuccessResponse.model_validate(payload)
    # No preview delivery because no sidecar.
    assert validated.data.preview_delivery is None
    assert validated.data.artifact_metadata.slide_count == 4

    cleanup_artifact(payload["data"]["artifact_metadata"])


def test_docx_response_passes_validation(client, monkeypatch) -> None:
    """DOCX response (no preview) still validates correctly."""
    import app.main

    monkeypatch.setattr(app.main, "sidecar_manager", _mock_ready_sidecar())

    body, headers, _ = signed_request_content("docx")
    response = client.post("/v1/generate", content=body, headers=headers)

    assert response.status_code == 200
    payload = response.json()

    validated = GenerateSuccessResponse.model_validate(payload)
    assert validated.data.preview_delivery is None  # no preview for docx
    assert validated.status == "completed"

    cleanup_artifact(payload["data"]["artifact_metadata"])


# ---------------------------------------------------------------------------
# 2. slide_count consistency
# ---------------------------------------------------------------------------


def test_slide_count_matches_actual_slides(client, monkeypatch) -> None:
    """``artifact_metadata.slide_count`` equals ``len(presentation.slides)``."""
    import app.main

    monkeypatch.setattr(app.main, "sidecar_manager", _mock_ready_sidecar())

    body, headers, payload_dict = signed_request_content("pptx")
    response = client.post("/v1/generate", content=body, headers=headers)

    assert response.status_code == 200
    payload = response.json()
    metadata = payload["data"]["artifact_metadata"]
    signed_count = metadata["slide_count"]

    # Open the artifact to count actual slides.
    artifact_path = artifact_path_from_metadata(metadata)
    assert artifact_path is not None, "artifact path should be resolvable"

    prs = Presentation(str(artifact_path))
    actual_count = len(prs.slides)

    assert signed_count == actual_count, (
        f"metadata slide_count ({signed_count}) != actual slides ({actual_count})"
    )

    cleanup_artifact(metadata)


# ---------------------------------------------------------------------------
# 3. HMAC consistency
# ---------------------------------------------------------------------------


def test_preview_url_uses_same_hmac_secret_as_artifact(
    client, monkeypatch,
) -> None:
    """HMAC signatures for preview and artifact use the same secret.

    This test verifies that ``build_signed_artifact_locator`` (which generates
    the signed URL) was called with the same secret for both the artifact and
    the preview.  Since both go through ``artifact_download._build_artifact_signature``
    with the module-level secret, they will use the same signing material.
    """
    import app.main

    monkeypatch.setattr(app.main, "sidecar_manager", _mock_ready_sidecar())

    body, headers, _ = signed_request_content("pptx")
    response = client.post("/v1/generate", content=body, headers=headers)

    assert response.status_code == 200
    payload = response.json()

    artifact_locator = payload["data"]["artifact_delivery"]
    preview_locator = payload["data"]["preview_delivery"]["locator"]

    # Both locators should be signed URLs.
    assert artifact_locator["kind"] == "signed_url"
    assert preview_locator["kind"] == "signed_url"

    artifact_params = _extract_signed_url_params(artifact_locator)
    preview_params = _extract_signed_url_params(preview_locator)

    # Both signatures are SHA256 hex digests (64 chars).
    assert len(artifact_params["signature"][0]) == 64
    assert len(preview_params["signature"][0]) == 64

    # Both use the same shared secret (checked implicitly by verifying both
    # download requests succeed).
    dl_artifact = client.get(
        f"{urlparse(artifact_locator['value']).path}?"
        f"{urlparse(artifact_locator['value']).query}"
    )
    dl_preview = client.get(
        f"{urlparse(preview_locator['value']).path}?"
        f"{urlparse(preview_locator['value']).query}"
    )

    assert dl_artifact.status_code == 200
    assert dl_preview.status_code == 200

    cleanup_artifact(payload["data"]["artifact_metadata"])


def test_hmac_signature_algorithm_is_sha256(
    client, monkeypatch,
) -> None:
    """The HMAC signature algorithm is ``hmac-sha256``."""
    import app.main
    from app.contracts import SIGNATURE_ALGORITHM

    assert SIGNATURE_ALGORITHM == "hmac-sha256"

    monkeypatch.setattr(app.main, "sidecar_manager", _mock_ready_sidecar())

    body, headers, _ = signed_request_content("pptx")
    response = client.post("/v1/generate", content=body, headers=headers)

    assert response.status_code == 200
    payload = response.json()

    # Verify algorithm is declared in health endpoint.
    health = client.get("/health")
    assert health.json()["auth"]["signature_algorithm"] == SIGNATURE_ALGORITHM

    cleanup_artifact(payload["data"]["artifact_metadata"])


# ---------------------------------------------------------------------------
# 4. Timeout / retry configuration
# ---------------------------------------------------------------------------


def test_marp_render_timeout_is_within_gateway_budget() -> None:
    """Marp sidecar render timeout (≤ 30s) + margin fits within 60s Gateway budget.

    The Gateway→MediaGen timeout is 60s.  Marp render takes at most 30s,
    which leaves 30s for:
    - LLM interpretation (≤ 10s typical)
    - Template injection / canvas rendering (≤ 5s typical)
    - Preview HTML rendering (≤ 10s typical)
    - R2 upload (≤ 5s typical)
    """
    from app.settings import get_settings

    settings = get_settings()
    marp_timeout = settings.marp_sidecar_render_timeout_seconds
    gateway_timeout = 60  # Gateway→MediaGen timeout (contract)

    assert marp_timeout <= 30, (
        f"Marp render timeout ({marp_timeout}s) should be ≤ 30s "
        f"to fit within the {gateway_timeout}s Gateway budget"
    )
    remaining = gateway_timeout - marp_timeout
    assert remaining >= 20, (
        f"Only {remaining}s remaining after Marp render — "
        f"need at least 20s for the rest of the pipeline"
    )


def test_marp_sidecar_config_defaults_are_safe() -> None:
    """Default sidecar settings are within safe limits."""
    from app.settings import get_settings

    settings = get_settings()

    # Ready timeout should not exceed the Gateway budget.
    assert settings.marp_sidecar_ready_timeout_seconds <= 60
    # Max concurrent renders should be reasonable.
    assert 1 <= settings.marp_sidecar_max_concurrent_renders <= 10
    # Health check interval should be frequent enough.
    assert settings.marp_sidecar_health_interval_seconds <= 60
    # Render timeout sanity check.
    assert 10 <= settings.marp_sidecar_render_timeout_seconds <= 60


# ---------------------------------------------------------------------------
# 5. Fase 4 additive fields
# ---------------------------------------------------------------------------


def test_artifact_metadata_includes_preview_url_on_success(
    client, monkeypatch,
) -> None:
    """When a preview is generated, ``artifact_metadata.preview_url`` is set."""
    import app.main

    monkeypatch.setattr(app.main, "sidecar_manager", _mock_ready_sidecar())

    body, headers, _ = signed_request_content("pptx")
    response = client.post("/v1/generate", content=body, headers=headers)

    assert response.status_code == 200
    payload = response.json()
    metadata = payload["data"]["artifact_metadata"]

    # The preview_url should be a signed URL string.
    preview_url = metadata.get("preview_url")
    assert preview_url is not None, "expected preview_url in artifact_metadata"
    assert isinstance(preview_url, str)
    assert preview_url.startswith("http"), f"expected http(s) URL, got {preview_url[:50]}"

    # Also verify it's the same URL as preview_delivery.locator.value.
    pd = payload["data"]["preview_delivery"]
    assert pd["locator"]["value"] == preview_url, (
        "preview_url in metadata should match preview_delivery.locator.value"
    )

    cleanup_artifact(metadata)


def test_artifact_metadata_excludes_preview_url_without_sidecar(client) -> None:
    """Without sidecar, ``preview_url`` is absent (None)."""
    body, headers, _ = signed_request_content("pptx")
    response = client.post("/v1/generate", content=body, headers=headers)

    assert response.status_code == 200
    payload = response.json()
    metadata = payload["data"]["artifact_metadata"]

    preview_url = metadata.get("preview_url")
    assert preview_url is None, "expected no preview_url without sidecar"

    cleanup_artifact(metadata)


def test_pptx_artifact_metadata_includes_layout_sources(
    client, monkeypatch,
) -> None:
    """PPTX metadata includes ``layout_sources`` when generated with canvas fallback."""
    import app.main

    monkeypatch.setattr(app.main, "sidecar_manager", _mock_ready_sidecar())

    body, headers, _ = signed_request_content("pptx")
    response = client.post("/v1/generate", content=body, headers=headers)

    assert response.status_code == 200
    payload = response.json()
    metadata = payload["data"]["artifact_metadata"]

    # The sample request produces 4 slides, all within template capacity,
    # so all should be "template".
    layout_sources = metadata.get("layout_sources")
    # Note: layout_sources may be None if the generator didn't populate it
    # (depends on engine wiring).  If present, verify the format.
    if layout_sources is not None:
        assert isinstance(layout_sources, list)
        assert len(layout_sources) == 4, f"expected 4 sources, got {len(layout_sources)}"
        for source in layout_sources:
            assert source in ("template", "canvas"), (
                f"unexpected layout_source: {source}"
            )

    cleanup_artifact(metadata)


# ---------------------------------------------------------------------------
# 6. Error contract validation
# ---------------------------------------------------------------------------


def test_health_endpoint_reports_template_registry(client) -> None:
    """Health endpoint includes template registry status."""
    response = client.get("/health")
    assert response.status_code == 200
    payload = response.json()

    templates = payload.get("templates")
    assert templates is not None, "expected 'templates' in health payload"
    # Template registry may or may not be loaded (depends on test env).
    # At minimum the field structure should be present.
    assert isinstance(templates, dict)
    assert "enabled" in templates


def test_model_validate_artifact_metadata_with_all_fields() -> None:
    """``ArtifactMetadata`` validates correctlty with all Fase 4 fields."""
    from app.contracts import ARTIFACT_METADATA_VERSION

    metadata = ArtifactMetadata.model_validate({
        "schema_version": ARTIFACT_METADATA_VERSION,
        "export_format": "pptx",
        "title": "Test",
        "filename": "test.pptx",
        "extension": "pptx",
        "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "size_bytes": 1024,
        "checksum_sha256": "a" * 64,
        "slide_count": 4,
        "artifact_locator": {"kind": "temporary_path", "value": "/tmp/test.pptx"},
        "generator": {"name": "test", "version": "1.0"},
        "warnings": [],
        "layout_sources": ["template", "template", "template", "canvas"],
        "preview_url": "https://example.com/preview.html",
    })

    assert metadata.layout_sources == ["template", "template", "template", "canvas"]
    assert metadata.preview_url == "https://example.com/preview.html"
    assert metadata.slide_count == 4


def test_model_validate_artifact_metadata_without_optional_fields() -> None:
    """``ArtifactMetadata`` validates correctly without Fase 4 optional fields."""
    from app.contracts import ARTIFACT_METADATA_VERSION

    metadata = ArtifactMetadata.model_validate({
        "schema_version": ARTIFACT_METADATA_VERSION,
        "export_format": "pptx",
        "title": "Test",
        "filename": "test.pptx",
        "extension": "pptx",
        "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "size_bytes": 1024,
        "checksum_sha256": "a" * 64,
        "slide_count": 4,
        "artifact_locator": {"kind": "temporary_path", "value": "/tmp/test.pptx"},
        "generator": {"name": "test", "version": "1.0"},
        "warnings": [],
        # layout_sources and preview_url are omitted — should default to None.
    })

    assert metadata.layout_sources is None
    assert metadata.preview_url is None

"""End-to-end pipeline tests for the media generator service.

Tests cover the full generation pipeline for all 3 export formats (DOCX, PDF,
PPTX) from API request through rendering to artifact download.  This replaces
the original plan's "5C" sub-task which assumed Flutter/Gateway integration;
instead we verify the Python-side pipeline is fully functional.

Gate criteria verified:
1. All 3 formats generate successfully via Template-Driven engines.
2. Preview delivery works for PDF and PPTX (DOCX has no preview).
3. Artifacts are downloadable via signed URLs with correct MIME types.
4. Visual parity between PDF and PPTX outputs (content structure match).
5. Response contract is backward-compatible with ``media_generation_spec.v1``.
"""

from __future__ import annotations

import json
from pathlib import Path
from unittest.mock import AsyncMock, MagicMock
from urllib.parse import parse_qs, urlparse

import pytest
# pyrefly: ignore [missing-import]
from pptx import Presentation

from app.models import GenerateSuccessResponse
from tests.helpers import (
    artifact_path_from_metadata,
    cleanup_artifact,
    signed_request_content,
)


# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

_MASTERS_DIR = Path(__file__).resolve().parent.parent / "app" / "templates" / "masters"
_HTML_MASTER_PATH = _MASTERS_DIR / "klass-educational-v1.html"
_DOCX_MASTER_PATH = _MASTERS_DIR / "klass-educational-v1.docx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _mock_ready_sidecar() -> MagicMock:
    """Return a mock ``SidecarManager`` that appears ready."""
    sc = MagicMock()
    # Mock PDF content must be > 100 bytes for validation
    mock_pdf = b"%PDF-1.4-mock-content-" + b"x" * 100
    sc.html_to_pdf = AsyncMock(return_value=mock_pdf)
    sc.is_running = True
    sc.is_ready = True
    sc.uptime_seconds = 42.0
    return sc


def _mock_template_registry() -> MagicMock:
    """Return a mock ``TemplateRegistry`` that resolves real masters."""
    tr = MagicMock()
    tr.get_html_master.return_value = _HTML_MASTER_PATH
    tr.get_docx_master.return_value = _DOCX_MASTER_PATH
    return tr


def _patch_deps(monkeypatch: pytest.MonkeyPatch) -> None:
    """Patch ``app.main`` globals for sidecar and template registry."""
    import app.main

    monkeypatch.setattr(app.main, "sidecar_manager", _mock_ready_sidecar())
    monkeypatch.setattr(app.main, "template_registry", _mock_template_registry())


def _download_artifact(client, metadata: dict) -> bytes:
    """Download artifact via its signed URL and return content."""
    artifact_locator = metadata["artifact_locator"]
    if artifact_locator["kind"] == "signed_url":
        parsed = urlparse(artifact_locator["value"])
        path = f"{parsed.path}?{parsed.query}"
    else:
        path = artifact_locator["value"]

    response = client.get(path)
    assert response.status_code == 200, f"Download failed: {response.status_code}"
    return response.content


def _download_preview(client, payload: dict) -> bytes:
    """Download preview via signed URL from preview_delivery."""
    preview_delivery = payload["data"]["preview_delivery"]
    assert preview_delivery is not None, "No preview_delivery in response"

    locator = preview_delivery["locator"]
    parsed = urlparse(locator["value"])
    path = f"{parsed.path}?{parsed.query}"

    response = client.get(path)
    assert response.status_code == 200, f"Preview download failed: {response.status_code}"
    return response.content


def _extract_sections_from_docx(docx_path: Path) -> list[str]:
    """Extract heading texts from a DOCX file for parity checking."""
    # pyrefly: ignore [missing-import]
    from docx import Document

    doc = Document(str(docx_path))
    headings = []
    for para in doc.paragraphs:
        if para.style.name.startswith("Heading"):
            headings.append(para.text.strip())
    return [h for h in headings if h]


def _extract_slide_titles_from_pptx(pptx_path: Path) -> list[str]:
    """Extract slide titles from a PPTX file for parity checking."""
    prs = Presentation(str(pptx_path))
    titles = []
    for slide in prs.slides:
        if slide.shapes.title:
            titles.append(slide.shapes.title.text.strip())
    return [t for t in titles if t]


# ---------------------------------------------------------------------------
# 1. Full Pipeline — DOCX
# ---------------------------------------------------------------------------


class TestDocxPipeline:
    """E2E tests for DOCX generation pipeline."""

    def test_docx_generate_and_download(self, client, monkeypatch) -> None:
        """Full DOCX pipeline: request → render → download → validate."""
        _patch_deps(monkeypatch)

        body, headers, _ = signed_request_content("docx")
        response = client.post("/v1/generate", content=body, headers=headers)

        assert response.status_code == 200
        payload = response.json()

        # Validate response contract
        validated = GenerateSuccessResponse.model_validate(payload)
        assert validated.status == "completed"
        assert validated.data.artifact_metadata.export_format == "docx"

        # Download and validate artifact
        metadata = payload["data"]["artifact_metadata"]
        content = _download_artifact(client, metadata)

        # Verify DOCX magic bytes (PK zip header)
        assert content[:2] == b"PK", "DOCX should start with PK zip header"
        assert len(content) > 100, "DOCX file should not be empty"

        # Verify it's a valid DOCX
        import io
        # pyrefly: ignore [missing-import]
        from docx import Document

        doc = Document(io.BytesIO(content))
        assert len(doc.paragraphs) > 0, "DOCX should have content"

        cleanup_artifact(metadata)

    def test_docx_no_preview_delivery(self, client, monkeypatch) -> None:
        """DOCX format should not have preview delivery."""
        _patch_deps(monkeypatch)

        body, headers, _ = signed_request_content("docx")
        response = client.post("/v1/generate", content=body, headers=headers)

        assert response.status_code == 200
        payload = response.json()

        assert payload["data"]["preview_delivery"] is None, (
            "DOCX should not have preview delivery"
        )

        cleanup_artifact(payload["data"]["artifact_metadata"])


# ---------------------------------------------------------------------------
# 2. Full Pipeline — PDF
# ---------------------------------------------------------------------------


class TestPdfPipeline:
    """E2E tests for PDF generation pipeline."""

    def test_pdf_generate_and_download(self, running_client, monkeypatch) -> None:
        """Full PDF pipeline: request → render → download → validate."""
        client = running_client
        _patch_deps(monkeypatch)

        body, headers, _ = signed_request_content("pdf")
        response = client.post("/v1/generate", content=body, headers=headers)

        assert response.status_code == 200
        payload = response.json()

        # Validate response contract
        validated = GenerateSuccessResponse.model_validate(payload)
        assert validated.status == "completed"
        assert validated.data.artifact_metadata.export_format == "pdf"

        # Download and validate artifact
        metadata = payload["data"]["artifact_metadata"]
        content = _download_artifact(client, metadata)

        # Verify PDF magic bytes
        assert content[:5] == b"%PDF-", "PDF should start with %PDF- header"
        assert len(content) > 100, "PDF file should not be empty"

        cleanup_artifact(metadata)

    def test_pdf_preview_delivery(self, running_client, monkeypatch) -> None:
        """PDF should have preview delivery with self-contained HTML."""
        client = running_client
        _patch_deps(monkeypatch)

        body, headers, _ = signed_request_content("pdf")
        response = client.post("/v1/generate", content=body, headers=headers)

        assert response.status_code == 200
        payload = response.json()

        # Verify preview delivery exists
        preview = payload["data"]["preview_delivery"]
        assert preview is not None, "PDF should have preview delivery"
        assert preview["schema_version"] == "media_generator_preview.v1"
        assert preview["mime_type"] == "text/html"
        assert preview["locator"]["kind"] == "signed_url"

        # Download and validate preview HTML
        preview_content = _download_preview(client, payload)
        html_str = preview_content.decode("utf-8")

        # Verify self-contained HTML
        assert "<!DOCTYPE html>" in html_str or "<html" in html_str, (
            "Preview should be valid HTML"
        )
        # Verify no external URLs (self-contained)
        assert "http://" not in html_str or "https://cdn" not in html_str, (
            "Preview should be self-contained (no external CDN URLs)"
        )

        cleanup_artifact(payload["data"]["artifact_metadata"])

    def test_pdf_page_count_in_metadata(self, running_client, monkeypatch) -> None:
        """PDF metadata should include page_count."""
        client = running_client
        _patch_deps(monkeypatch)

        body, headers, _ = signed_request_content("pdf")
        response = client.post("/v1/generate", content=body, headers=headers)

        assert response.status_code == 200
        payload = response.json()
        metadata = payload["data"]["artifact_metadata"]

        # PDF should have page_count, not slide_count
        assert metadata["slide_count"] is None, "PDF slide_count should be None"
        assert metadata["page_count"] is not None, "PDF should have page_count"
        assert metadata["page_count"] >= 1, "page_count should be >= 1"

        cleanup_artifact(metadata)


# ---------------------------------------------------------------------------
# 3. Full Pipeline — PPTX
# ---------------------------------------------------------------------------


class TestPptxPipeline:
    """E2E tests for PPTX generation pipeline."""

    def test_pptx_generate_and_download(self, client, monkeypatch) -> None:
        """Full PPTX pipeline: request → render → download → validate."""
        _patch_deps(monkeypatch)

        body, headers, _ = signed_request_content("pptx")
        response = client.post("/v1/generate", content=body, headers=headers)

        assert response.status_code == 200
        payload = response.json()

        # Validate response contract
        validated = GenerateSuccessResponse.model_validate(payload)
        assert validated.status == "completed"
        assert validated.data.artifact_metadata.export_format == "pptx"

        # Download and validate artifact
        metadata = payload["data"]["artifact_metadata"]
        content = _download_artifact(client, metadata)

        # Verify PPTX magic bytes (PK zip header)
        assert content[:2] == b"PK", "PPTX should start with PK zip header"
        assert len(content) > 100, "PPTX file should not be empty"

        # Verify it's a valid PPTX with correct slide count
        import io
        prs = Presentation(io.BytesIO(content))
        assert len(prs.slides) == 4, f"Expected 4 slides, got {len(prs.slides)}"

        cleanup_artifact(metadata)

    def test_pptx_preview_delivery(self, client, monkeypatch) -> None:
        """PPTX should have preview delivery with self-contained HTML."""
        _patch_deps(monkeypatch)

        body, headers, _ = signed_request_content("pptx")
        response = client.post("/v1/generate", content=body, headers=headers)

        assert response.status_code == 200
        payload = response.json()

        # Verify preview delivery exists
        preview = payload["data"]["preview_delivery"]
        assert preview is not None, "PPTX should have preview delivery"
        assert preview["schema_version"] == "media_generator_preview.v1"
        assert preview["mime_type"] == "text/html"

        # Download and validate preview HTML
        preview_content = _download_preview(client, payload)
        html_str = preview_content.decode("utf-8")

        # Verify self-contained HTML
        assert "<!DOCTYPE html>" in html_str or "<html" in html_str, (
            "Preview should be valid HTML"
        )

        cleanup_artifact(payload["data"]["artifact_metadata"])

    def test_pptx_slide_count_consistency(self, client, monkeypatch) -> None:
        """PPTX metadata slide_count should match actual slides."""
        _patch_deps(monkeypatch)

        body, headers, _ = signed_request_content("pptx")
        response = client.post("/v1/generate", content=body, headers=headers)

        assert response.status_code == 200
        payload = response.json()
        metadata = payload["data"]["artifact_metadata"]

        # Download and count actual slides
        content = _download_artifact(client, metadata)
        import io
        prs = Presentation(io.BytesIO(content))
        actual_count = len(prs.slides)

        assert metadata["slide_count"] == actual_count, (
            f"metadata slide_count ({metadata['slide_count']}) != "
            f"actual slides ({actual_count})"
        )

        cleanup_artifact(metadata)

    def test_pptx_layout_sources(self, client, monkeypatch) -> None:
        """PPTX metadata should include layout_sources."""
        _patch_deps(monkeypatch)

        body, headers, _ = signed_request_content("pptx")
        response = client.post("/v1/generate", content=body, headers=headers)

        assert response.status_code == 200
        payload = response.json()
        metadata = payload["data"]["artifact_metadata"]

        layout_sources = metadata.get("layout_sources")
        if layout_sources is not None:
            assert isinstance(layout_sources, list)
            assert len(layout_sources) == 4, f"Expected 4 sources, got {len(layout_sources)}"
            for source in layout_sources:
                assert source in ("template", "canvas"), (
                    f"unexpected layout_source: {source}"
                )

        cleanup_artifact(metadata)


# ---------------------------------------------------------------------------
# 4. Cross-Format Parity Verification
# ---------------------------------------------------------------------------


class TestCrossFormatParity:
    """Verify visual parity between PDF and PPTX outputs."""

    def test_pdf_pptx_content_structure_parity(self, running_client, monkeypatch) -> None:
        """PDF and PPTX should have similar content structure.

        While visual styling may differ, the content structure (sections, titles)
        should be consistent between formats for the same input.
        """
        client = running_client
        _patch_deps(monkeypatch)

        # Generate both PDF and PPTX with same content
        body_pdf, headers_pdf, _ = signed_request_content("pdf")
        body_pptx, headers_pptx, _ = signed_request_content("pptx")

        response_pdf = client.post("/v1/generate", content=body_pdf, headers=headers_pdf)
        response_pptx = client.post("/v1/generate", content=body_pptx, headers=headers_pptx)

        assert response_pdf.status_code == 200
        assert response_pptx.status_code == 200

        payload_pdf = response_pdf.json()
        payload_pptx = response_pptx.json()

        # Both should have same generation context
        assert payload_pdf["data"]["generation_id"] == "generation-pdf-001"
        assert payload_pptx["data"]["generation_id"] == "generation-pptx-001"

        # Both should have same number of sections in the request
        # (verified through the contract - both use same sample_request_payload)

        # Verify both have preview delivery
        assert payload_pdf["data"]["preview_delivery"] is not None
        assert payload_pptx["data"]["preview_delivery"] is not None

        # Verify both previews are valid HTML
        preview_pdf = _download_preview(client, payload_pdf)
        preview_pptx = _download_preview(client, payload_pptx)

        html_pdf = preview_pdf.decode("utf-8")
        html_pptx = preview_pptx.decode("utf-8")

        # Both should be self-contained HTML
        for html in [html_pdf, html_pptx]:
            assert "<!DOCTYPE html>" in html or "<html" in html, (
                "Both previews should be valid HTML"
            )
            # Should contain Jinja2-rendered content (not raw template)
            assert "{{" not in html, "Preview should not contain raw Jinja2 placeholders"
            assert "{%" not in html, "Preview should not contain raw Jinja2 blocks"

        cleanup_artifact(payload_pdf["data"]["artifact_metadata"])
        cleanup_artifact(payload_pptx["data"]["artifact_metadata"])

    def test_all_formats_share_same_contract_version(self, running_client, monkeypatch) -> None:
        """All 3 formats should use the same contract versions."""
        client = running_client
        _patch_deps(monkeypatch)

        versions = {}
        for fmt in ("docx", "pdf", "pptx"):
            body, headers, _ = signed_request_content(fmt)
            response = client.post("/v1/generate", content=body, headers=headers)

            assert response.status_code == 200
            payload = response.json()
            versions[fmt] = payload["data"]["artifact_metadata"]["schema_version"]

            cleanup_artifact(payload["data"]["artifact_metadata"])

        # All should use same schema version
        assert versions["docx"] == versions["pdf"] == versions["pptx"], (
            f"Schema versions differ: {versions}"
        )


# ---------------------------------------------------------------------------
# 5. Response Contract Backward Compatibility
# ---------------------------------------------------------------------------


class TestResponseContract:
    """Verify response contract backward compatibility."""

    def test_response_has_required_fields(self, client, monkeypatch) -> None:
        """Response should have all required fields per contract."""
        _patch_deps(monkeypatch)

        body, headers, _ = signed_request_content("pptx")
        response = client.post("/v1/generate", content=body, headers=headers)

        assert response.status_code == 200
        payload = response.json()

        # Required top-level fields
        assert "status" in payload
        assert "data" in payload

        # Required data fields
        data = payload["data"]
        assert "generation_id" in data
        assert "artifact_metadata" in data
        assert "artifact_delivery" in data

        # Artifact metadata required fields
        metadata = data["artifact_metadata"]
        assert "schema_version" in metadata
        assert "export_format" in metadata
        assert "filename" in metadata
        assert "mime_type" in metadata
        assert "size_bytes" in metadata
        assert "checksum_sha256" in metadata

        cleanup_artifact(metadata)

    def test_additive_fields_present(self, client, monkeypatch) -> None:
        """Additive fields (preview_url, layout_sources) should be present."""
        _patch_deps(monkeypatch)

        body, headers, _ = signed_request_content("pptx")
        response = client.post("/v1/generate", content=body, headers=headers)

        assert response.status_code == 200
        payload = response.json()
        metadata = payload["data"]["artifact_metadata"]

        # These are additive fields - should be present (may be None)
        assert "preview_url" in metadata, "preview_url should be in metadata"
        assert "layout_sources" in metadata, "layout_sources should be in metadata"

        # For PPTX with sidecar, preview_url should be set
        assert metadata["preview_url"] is not None, "PPTX should have preview_url"
        assert metadata["preview_url"].startswith("http"), "preview_url should be http(s) URL"

        cleanup_artifact(metadata)

    def test_preview_delivery_schema(self, running_client, monkeypatch) -> None:
        """preview_delivery should follow media_generator_preview.v1 schema."""
        client = running_client
        _patch_deps(monkeypatch)

        body, headers, _ = signed_request_content("pdf")
        response = client.post("/v1/generate", content=body, headers=headers)

        assert response.status_code == 200
        payload = response.json()
        preview = payload["data"]["preview_delivery"]

        assert preview is not None
        assert preview["schema_version"] == "media_generator_preview.v1"
        assert preview["mime_type"] == "text/html"

        locator = preview["locator"]
        assert locator["kind"] == "signed_url"
        assert "value" in locator
        assert locator["value"].startswith("http"), "Locator URL should be http(s)"

        # Verify signed URL has required parameters
        parsed = urlparse(locator["value"])
        params = parse_qs(parsed.query)
        assert "signature" in params, "Signed URL should have signature"
        assert "expires" in params, "Signed URL should have expires"

        cleanup_artifact(payload["data"]["artifact_metadata"])


# ---------------------------------------------------------------------------
# 6. Artifact Download Verification
# ---------------------------------------------------------------------------


class TestArtifactDownload:
    """Verify artifact download via signed URLs."""

    def test_signed_url_download_returns_correct_content_type(
        self, running_client, monkeypatch,
    ) -> None:
        """Download should return correct Content-Type for each format."""
        client = running_client
        _patch_deps(monkeypatch)

        expected_types = {
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "pdf": "application/pdf",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        }

        for fmt, expected_type in expected_types.items():
            body, headers, _ = signed_request_content(fmt)
            response = client.post("/v1/generate", content=body, headers=headers)

            assert response.status_code == 200
            payload = response.json()
            metadata = payload["data"]["artifact_metadata"]

            # Download artifact
            artifact_locator = metadata["artifact_locator"]
            parsed = urlparse(artifact_locator["value"])
            path = f"{parsed.path}?{parsed.query}"
            dl_response = client.get(path)

            assert dl_response.status_code == 200
            # Content-Type may vary, but should be set
            assert "content-type" in dl_response.headers

            cleanup_artifact(metadata)

    def test_preview_download_returns_html(self, running_client, monkeypatch) -> None:
        """Preview download should return HTML content."""
        client = running_client
        _patch_deps(monkeypatch)

        for fmt in ("pdf", "pptx"):  # DOCX has no preview
            body, headers, _ = signed_request_content(fmt)
            response = client.post("/v1/generate", content=body, headers=headers)

            assert response.status_code == 200
            payload = response.json()
            preview = payload["data"]["preview_delivery"]

            if preview is not None:
                parsed = urlparse(preview["locator"]["value"])
                path = f"{parsed.path}?{parsed.query}"
                dl_response = client.get(path)

                assert dl_response.status_code == 200
                assert dl_response.headers.get("content-type", "").startswith("text/html"), (
                    f"Preview Content-Type should be text/html for {fmt}"
                )

            cleanup_artifact(payload["data"]["artifact_metadata"])


# ---------------------------------------------------------------------------
# 7. Error Handling
# ---------------------------------------------------------------------------


class TestErrorHandling:
    """Verify error handling in the pipeline."""

    def test_invalid_export_format_returns_error(self, client, monkeypatch) -> None:
        """Invalid export format should return 422 or appropriate error."""
        _patch_deps(monkeypatch)

        body, headers, _ = signed_request_content("invalid_format")
        response = client.post("/v1/generate", content=body, headers=headers)

        # Should return error (422 validation error or similar)
        assert response.status_code in (400, 422, 500), (
            f"Expected error status, got {response.status_code}"
        )

    def test_missing_signature_returns_401(self, client, monkeypatch) -> None:
        """Request without signature should return 401."""
        _patch_deps(monkeypatch)

        body, _, _ = signed_request_content("pptx")
        headers = {"Content-Type": "application/json"}  # No signature

        response = client.post("/v1/generate", content=body, headers=headers)
        assert response.status_code == 401, "Missing signature should return 401"

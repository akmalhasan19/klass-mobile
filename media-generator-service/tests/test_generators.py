from __future__ import annotations

from pathlib import Path

from docx import Document
from pptx import Presentation

from app.document_model import build_render_document
from app.generators.docx_generator import DocxGenerator
from app.generators.pdf_generator import PdfGenerator
from app.generators.pptx_generator import PptxGenerator
from app.generators.registry import GeneratorRegistry
from app.settings import get_settings
from tests.helpers import cleanup_artifact, sample_request


def test_docx_generator_renders_expected_sections_and_text() -> None:
    request_payload = sample_request("docx")
    render_document = build_render_document(request_payload.generation_spec)
    generator = DocxGenerator()

    metadata = generator.generate(request_payload, render_document, get_settings())
    document = Document(metadata["artifact_locator"]["value"])
    paragraph_text = "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())

    assert metadata["title"] == "Handout Pecahan Kelas 5"
    assert "Tujuan Belajar" in paragraph_text
    assert "Contoh dan Latihan" in paragraph_text
    assert "Latihan Mandiri" in paragraph_text

    cleanup_artifact(metadata)


def test_pdf_generator_writes_pdf_header_and_page_count() -> None:
    request_payload = sample_request("pdf")
    render_document = build_render_document(request_payload.generation_spec)
    generator = PdfGenerator()

    metadata = generator.generate(request_payload, render_document, get_settings())
    artifact_path = Path(metadata["artifact_locator"]["value"])

    assert artifact_path.read_bytes().startswith(b"%PDF")
    assert metadata["page_count"] >= 1
    assert metadata["mime_type"] == "application/pdf"

    cleanup_artifact(metadata)


def test_registry_only_exposes_docx_and_pdf_generators() -> None:
    registry = GeneratorRegistry()

    assert registry.get("docx").export_format == "docx"
    assert registry.get("pdf").export_format == "pdf"
    assert registry.get("pptx").export_format == "pptx"


def test_pptx_generator_renders_title_section_and_activity_slides() -> None:
    request_payload = sample_request("pptx")
    render_document = build_render_document(request_payload.generation_spec)
    generator = PptxGenerator()

    metadata = generator.generate(request_payload, render_document, get_settings())
    presentation = Presentation(metadata["artifact_locator"]["value"])
    slide_texts = []

    for slide in presentation.slides:
        fragments = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = getattr(shape, "text", "")
                if isinstance(text, str) and text.strip():
                    fragments.append(text.strip())
        slide_texts.append("\n".join(fragments))

    assert metadata["slide_count"] == 4
    assert any("Handout Pecahan Kelas 5" in text for text in slide_texts)
    assert any("Tujuan Belajar" in text for text in slide_texts)
    assert any("Aktivitas dan Penilaian" in text for text in slide_texts)

    cleanup_artifact(metadata)
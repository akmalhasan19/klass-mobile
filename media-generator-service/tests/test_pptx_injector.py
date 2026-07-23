"""Integration tests for the Template Injector (Fase 3).

Exercises the full injection pipeline: loading a real master template,
building a ``SlideBlueprint``, running ``TemplateInjector.inject()``,
and verifying the output ``.pptx`` has correct content, slide counts,
and is valid (re-openable by ``python-pptx``).

Coverage
--------
* Basic placeholder injection — title, subtitle, body text appear in shapes.
* Capacity overflow — slides exceeding ``max_cards`` / ``max_chars`` trigger
  the canvas fallback engine so the deck contains a mix of template-injected
  and canvas-calculated slides.
* Missing layout — slides whose ``slide_type`` has no manifest entry are
  recorded in ``fallback_slides`` and ``warnings``.
* Slide count accuracy — ``InjectionResult.slide_count`` matches the actual
  number of slides in the output file.
* File integrity — the output ``.pptx`` is re-openable with ``python-pptx``
  and has the expected slide dimensions (matching the master).
* Hybrid deck — a single output file contains both template-filled slides
  and canvas-calculated slides, each with distinct shape characteristics.
"""
from __future__ import annotations

from pathlib import Path

import pytest
from pydantic import ValidationError
# pyrefly: ignore [missing-import]
from pptx import Presentation

from app.engines.blueprint import (
    Card,
    ContentBlock,
    DeckMeta,
    Slide,
    SlideBlueprint,
)
from app.engines.canvas_calculator.layout_engine import CanvasLayoutEngine
from app.engines.canvas_calculator.shape_renderer import CanvasShapeRenderer
from app.engines.pptx_injector.injector import TemplateInjector
from app.templates.registry import TemplateRegistry

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

_TEMPLATES_DIR = Path(__file__).resolve().parent.parent / "app" / "templates"


def _load_registry() -> TemplateRegistry:
    """Load the single bundled template into a fresh registry."""
    registry = TemplateRegistry()
    registry.load_templates(_TEMPLATES_DIR)
    return registry


def _injector() -> tuple[TemplateInjector, TemplateRegistry]:
    """Return a ``(injector, registry)`` pair using the bundled master.

    The injector is created **without** a canvas engine by default so callers
    can explicitly opt in when testing overflow / hybrid scenarios.
    """
    registry = _load_registry()
    entry = registry.get("klass-educational-v1")
    injector = TemplateInjector(
        master_path=entry.master_path,
        manifest=entry.manifest,
        canvas_engine=None,
    )
    return injector, registry


def _card(text: str, kind: str = "paragraph") -> Card:
    return Card(body_blocks=[ContentBlock(kind=kind, content=text)])


def _minimal_slide(slide_type: str, title: str = "Judul") -> Slide:
    return Slide(slide_type=slide_type, title=title, cards=[_card("konten")])


def _small_blueprint() -> SlideBlueprint:
    """A blueprint whose slides all fit within manifest capacities.

    Title: 1 card (2 objectives) → max_cards=4  ✅
    Content: 1 card with 2 blocks → max_cards=6 ✅
    Assessment: 1 card → max_cards=8 ✅
    """
    return SlideBlueprint(
        deck_meta=DeckMeta(
            title="Test Deck",
            summary="Ringkasan.",
            language="id",
            audience_level="elementary",
            tone="encouraging",
        ),
        slides=[
            Slide(
                slide_type="title",
                title="Judul Presentasi",
                subtitle="Subjudul",
                cards=[_card("Learning objective 1", "bullet"), _card("LO 2", "bullet")],
            ),
            Slide(
                slide_type="content",
                title="Konten Slide",
                cards=[_card("Isi konten utama.", "paragraph")],
            ),
            Slide(
                slide_type="assessment",
                title="Aktivitas",
                cards=[Card(heading="Latihan", body_blocks=[ContentBlock(kind="paragraph", content="Kerjakan soal.")])],
            ),
        ],
    )


# ---------------------------------------------------------------------------
# 1. Basic placeholder injection
# ---------------------------------------------------------------------------


def test_injector_fills_title_placeholder() -> None:
    """Title text is written into the master's title shape."""
    injector, _ = _injector()
    out = _tmp_output()

    blueprint = _small_blueprint()
    result = injector.inject(blueprint, out)
    prs = Presentation(str(out))

    title_slide = prs.slides[0]
    shape_texts = {s.name: s.text for s in title_slide.shapes if hasattr(s, "text")}

    assert "Title 1" in shape_texts, "master should have a 'Title 1' shape"
    assert "Judul Presentasi" in shape_texts["Title 1"]
    assert result.slide_count >= 1


def test_injector_fills_subtitle_placeholder() -> None:
    """Subtitle text appears in the master's subtitle shape."""
    injector, _ = _injector()
    out = _tmp_output()

    blueprint = _small_blueprint()
    injector.inject(blueprint, out)
    prs = Presentation(str(out))

    title_slide = prs.slides[0]
    shape_texts = {s.name: s.text for s in title_slide.shapes if hasattr(s, "text")}

    assert "Subtitle 2" in shape_texts, "master should have a 'Subtitle 2' shape"
    assert "Subjudul" in shape_texts.get("Subtitle 2", "")


def test_injector_fills_body_placeholder() -> None:
    """Body text (cards + bullets) appears in the master's body shape."""
    injector, _ = _injector()
    out = _tmp_output()

    blueprint = _small_blueprint()
    injector.inject(blueprint, out)
    prs = Presentation(str(out))

    # Content slide (index 1 in blueprint → index 1 in output)
    content_slide = prs.slides[1]
    shape_texts = {s.name: s.text for s in content_slide.shapes if hasattr(s, "text")}

    body_text = shape_texts.get("Content Placeholder 3", "")
    assert "Isi konten utama." in body_text, f"body should contain card text: {body_text!r}"


def test_injector_handles_notes() -> None:
    """Speaker notes are written to the slide's notes pane."""
    injector, _ = _injector()
    out = _tmp_output()

    blueprint = SlideBlueprint(
        deck_meta=DeckMeta(
            title="Test", summary=".", language="id",
            audience_level="elementary", tone="n",
        ),
        slides=[
            Slide(
                slide_type="content",
                title="Slide A",
                cards=[_card("x")],
                speaker_notes="Ingatkan siswa tentang konsep dasar.",
            ),
        ],
    )
    injector.inject(blueprint, out)
    prs = Presentation(str(out))

    notes_slide = prs.slides[0].notes_slide
    assert notes_slide is not None
    assert "Ingatkan siswa" in notes_slide.notes_text_frame.text


# ---------------------------------------------------------------------------
# 2. Slide count & warnings
# ---------------------------------------------------------------------------


def test_slide_count_matches_blueprint() -> None:
    """``InjectionResult.slide_count`` equals the number of slides written."""
    injector, _ = _injector()
    out = _tmp_output()

    blueprint = _small_blueprint()
    result = injector.inject(blueprint, out)

    prs = Presentation(str(out))
    assert result.slide_count == len(prs.slides)
    assert result.slide_count == len(blueprint.slides)


def test_warnings_collected_for_missing_layout() -> None:
    """Slides with no matching layout produce warnings + fallback."""
    injector, _ = _injector()
    out = _tmp_output()

    blueprint = SlideBlueprint(
        deck_meta=DeckMeta(
            title="T", summary=".", language="id",
            audience_level="elementary", tone="n",
        ),
        slides=[
            _minimal_slide("section", "Section Slide"),  # no manifest layout for 'section'
        ],
    )
    result = injector.inject(blueprint, out)

    assert len(result.warnings) >= 1
    assert any("no layout" in w.lower() for w in result.warnings)
    assert 1 in result.fallback_slides  # the section slide


# ---------------------------------------------------------------------------
# 3. Canvas fallback on overflow (capacity gate)
# ---------------------------------------------------------------------------


def test_canvas_engine_is_called_on_overflow() -> None:
    """Slides exceeding max_cards are delegated to the canvas engine."""
    registry = _load_registry()
    entry = registry.get("klass-educational-v1")

    canvas_engine = CanvasLayoutEngine(
        slide_width=Presentation(str(entry.master_path)).slide_width,
        slide_height=Presentation(str(entry.master_path)).slide_height,
        renderer=CanvasShapeRenderer(),
    )
    injector = TemplateInjector(
        master_path=entry.master_path,
        manifest=entry.manifest,
        canvas_engine=canvas_engine,
    )
    out = _tmp_output()

    # Content layout has max_cards=6 — use 8 cards to overflow.
    blueprint = SlideBlueprint(
        deck_meta=DeckMeta(
            title="T", summary=".", language="id",
            audience_level="elementary", tone="n",
        ),
        slides=[
            Slide(
                slide_type="content",
                title="Overflow Slide",
                cards=[_card(f"Card {i}") for i in range(8)],
            ),
        ],
    )
    result = injector.inject(blueprint, out)

    # Fallback should be recorded
    assert any("exceeds" in w.lower() for w in result.warnings), (
        f"expected overflow warnings, got: {result.warnings}"
    )
    assert result.fallback_slides == [1]


def test_canvas_fallback_when_no_layout() -> None:
    """When no layout exists for a slide_type, canvas engine renders it."""
    registry = _load_registry()
    entry = registry.get("klass-educational-v1")

    canvas_engine = CanvasLayoutEngine(renderer=CanvasShapeRenderer())
    injector = TemplateInjector(
        master_path=entry.master_path,
        manifest=entry.manifest,
        canvas_engine=canvas_engine,
    )
    out = _tmp_output()

    blueprint = SlideBlueprint(
        deck_meta=DeckMeta(
            title="T", summary=".", language="id",
            audience_level="elementary", tone="n",
        ),
        slides=[
            _minimal_slide("section"),  # no 'section' layout in manifest
        ],
    )
    result = injector.inject(blueprint, out)
    prs = Presentation(str(out))

    assert len(prs.slides) == 1
    assert 1 in result.fallback_slides


# ---------------------------------------------------------------------------
# 4. Hybrid deck (template + canvas in one file)
# ---------------------------------------------------------------------------


def test_hybrid_deck_contains_both_template_and_canvas_slides() -> None:
    """A deck with some fitting and some overflowing slides produces both."""
    registry = _load_registry()
    entry = registry.get("klass-educational-v1")

    canvas_engine = CanvasLayoutEngine(renderer=CanvasShapeRenderer())
    injector = TemplateInjector(
        master_path=entry.master_path,
        manifest=entry.manifest,
        canvas_engine=canvas_engine,
    )
    out = _tmp_output()

    # Slide 1: title — fits (1 card ≤ 4)
    # Slide 2: content — fits (1 card ≤ 6)
    # Slide 3: content — overflows (10 cards > 6)
    blueprint = SlideBlueprint(
        deck_meta=DeckMeta(
            title="Hybrid Deck", summary=".", language="id",
            audience_level="elementary", tone="n",
        ),
        slides=[
            Slide(slide_type="title", title="Judul", cards=[_card("LO", "bullet")]),
            Slide(slide_type="content", title="Konten Biasa", cards=[_card("Isi")]),
            Slide(
                slide_type="content",
                title="Konten Overflow",
                cards=[_card(f"Card {i}") for i in range(10)],
            ),
        ],
    )
    result = injector.inject(blueprint, out)

    # 3 slides expected (all 3 in blueprint are rendered)
    assert result.slide_count == 3
    # Slide 3 (index 3, 1-based) should be fallback
    assert 3 in result.fallback_slides, f"expected slide 3 in fallback, got {result.fallback_slides}"
    assert len(result.warnings) >= 1


# ---------------------------------------------------------------------------
# 5. File integrity (valid, re-openable)
# ---------------------------------------------------------------------------


def test_output_file_is_valid_pptx() -> None:
    """The output ``.pptx`` is a valid ZIP-based PPTX that re-opens cleanly."""
    injector, _ = _injector()
    out = _tmp_output()

    injector.inject(_small_blueprint(), out)

    # Re-open with python-pptx — should not raise.
    prs = Presentation(str(out))
    assert len(prs.slides) >= 1


def test_output_retains_master_dimensions() -> None:
    """Output slide dimensions match the master template."""
    registry = _load_registry()
    entry = registry.get("klass-educational-v1")

    master_prs = Presentation(str(entry.master_path))
    injector = TemplateInjector(
        master_path=entry.master_path,
        manifest=entry.manifest,
    )
    out = _tmp_output()

    injector.inject(_small_blueprint(), out)
    out_prs = Presentation(str(out))

    assert out_prs.slide_width == master_prs.slide_width
    assert out_prs.slide_height == master_prs.slide_height


# ---------------------------------------------------------------------------
# 6. Edge cases
# ---------------------------------------------------------------------------


def test_empty_blueprint_raises_validation_error() -> None:
    """A blueprint must have at least one slide (Pydantic enforces this)."""
    with pytest.raises(ValidationError):
        SlideBlueprint(
            deck_meta=DeckMeta(
                title="T", summary=".", language="id",
                audience_level="elementary", tone="n",
            ),
            slides=[],
        )


def test_injector_without_canvas_engine_still_records_fallback() -> None:
    """Without a canvas engine, overflow slides are still counted as fallback."""
    injector, _ = _injector()  # no canvas engine
    out = _tmp_output()

    blueprint = SlideBlueprint(
        deck_meta=DeckMeta(
            title="T", summary=".", language="id",
            audience_level="elementary", tone="n",
        ),
        slides=[
            Slide(
                slide_type="section",  # no layout for 'section'
                title="Section",
                cards=[_card("x")],
            ),
        ],
    )
    result = injector.inject(blueprint, out)

    assert 1 in result.fallback_slides
    assert len(result.warnings) >= 1


# ---------------------------------------------------------------------------
# Helper
# ---------------------------------------------------------------------------

def _tmp_output() -> Path:
    import tempfile
    import os

    fd, path = tempfile.mkstemp(suffix=".pptx", prefix="test_injector_")
    os.close(fd)
    return Path(path)
